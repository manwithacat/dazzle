"""
PPTX pitch deck generator.

Generates professional investor pitch decks from PitchContext.
Uses python-pptx for slide creation with a dark navy theme.

python-pptx is an optional dependency. The generator checks for
its availability at runtime and returns a clear error if missing.
"""

from __future__ import annotations

import logging
from collections.abc import Callable
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

from dazzle.pitch.extractor import PitchContext
from dazzle.pitch.ir import BrandColors, ExtraSlide, ExtraSlideLayout

logger = logging.getLogger(__name__)


@dataclass
class GeneratorResult:
    """Result of a generator run."""

    success: bool
    output_path: Path | None = None
    files_created: list[str] = field(default_factory=list)
    error: str | None = None
    slide_count: int = 0


def _check_pptx_available() -> bool:
    """Check if python-pptx is available."""
    try:
        import pptx  # noqa: F401

        return True
    except ImportError:
        return False


def _resolve_colors(brand: BrandColors) -> dict[str, Any]:
    """Convert brand hex colors to RGBColor objects."""
    from pptx.dml.color import RGBColor

    def hex_to_rgb(hex_color: str) -> RGBColor:
        h = hex_color.lstrip("#")
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

    return {
        "primary": hex_to_rgb(brand.primary),
        "accent": hex_to_rgb(brand.accent),
        "highlight": hex_to_rgb(brand.highlight),
        "success": hex_to_rgb(brand.success),
        "light": hex_to_rgb(brand.light),
        "white": RGBColor(0xFF, 0xFF, 0xFF),
        "dark_text": hex_to_rgb(brand.primary),
        "muted": RGBColor(0x99, 0x99, 0x99),
        "font_family": brand.font_family,
    }


def _add_text_box(
    slide: Any,
    left: Any,
    top: Any,
    width: Any,
    height: Any,
    text: str,
    *,
    font_size: int = 18,
    bold: bool = False,
    color: Any = None,
    alignment: int | None = None,
    font_name: str | None = None,
) -> Any:
    """Add a text box to a slide."""
    from pptx.util import Pt

    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    if color:
        p.font.color.rgb = color
    if alignment is not None:
        p.alignment = alignment
    if font_name:
        p.font.name = font_name
    return txbox


def _create_dark_slide(prs: Any, colors: dict[str, Any]) -> Any:
    """Create a slide with dark background."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = colors["primary"]
    return slide


def _add_speaker_notes(slide: Any, text: str) -> None:
    """Add speaker notes to a slide."""
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text


def _fmt_currency(amount: int, currency: str = "GBP") -> str:
    """Format an amount as currency."""
    symbols = {"GBP": "\u00a3", "USD": "$", "EUR": "\u20ac"}
    symbol = symbols.get(currency, currency + " ")
    if amount >= 1_000_000_000:
        return f"{symbol}{amount / 1_000_000_000:.1f}B"
    elif amount >= 1_000_000:
        return f"{symbol}{amount / 1_000_000:.1f}M"
    elif amount >= 1_000:
        return f"{symbol}{amount / 1_000:.0f}K"
    return f"{symbol}{amount:,}"


def _add_card(
    slide: Any,
    left: Any,
    top: Any,
    width: Any,
    height: Any,
    *,
    fill_color: Any = None,
    border_color: Any = None,
) -> Any:
    """Add a rounded rectangle card shape to a slide."""
    from pptx.util import Pt

    shape = slide.shapes.add_shape(
        5,  # MSO_SHAPE.ROUNDED_RECTANGLE
        left,
        top,
        width,
        height,
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def _add_stat_box(
    slide: Any,
    left: Any,
    top: Any,
    width: Any,
    value: str,
    label: str,
    *,
    value_color: Any = None,
    label_color: Any = None,
    alignment: int | None = None,
) -> None:
    """Add a big-number + caption stat box."""
    from pptx.util import Inches

    _add_text_box(
        slide,
        left,
        top,
        width,
        Inches(1),
        value,
        font_size=48,
        bold=True,
        color=value_color,
        alignment=alignment,
    )
    _add_text_box(
        slide,
        left,
        top + Inches(1),
        width,
        Inches(0.5),
        label,
        font_size=16,
        color=label_color,
        alignment=alignment,
    )


def _add_columns(
    slide: Any,
    top: Any,
    items: list[tuple[str, str]],
    *,
    value_color: Any = None,
    label_color: Any = None,
    alignment: int | None = None,
) -> None:
    """Auto-layout N stat boxes across slide width."""
    from pptx.util import Inches

    n = len(items)
    if n == 0:
        return
    box_w = min(3.0, 11.0 / max(n, 1))
    x_start = 1.0
    for i, (value, label) in enumerate(items[:6]):
        x = x_start + i * (box_w + 0.2)
        _add_stat_box(
            slide,
            Inches(x),
            top,
            Inches(box_w),
            value,
            label,
            value_color=value_color,
            label_color=label_color,
            alignment=alignment,
        )


def _add_rich_text_box(
    slide: Any,
    left: Any,
    top: Any,
    width: Any,
    height: Any,
    text: str,
    *,
    font_size: int = 18,
    color: Any = None,
    alignment: int | None = None,
    font_name: str | None = None,
) -> Any:
    """Add a text box that parses **bold** markers into bold runs."""
    import re

    from pptx.util import Pt

    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if alignment is not None:
        p.alignment = alignment

    parts = re.split(r"(\*\*.*?\*\*)", text)
    for part in parts:
        if part.startswith("**") and part.endswith("**"):
            run = p.add_run()
            run.text = part[2:-2]
            run.font.bold = True
            run.font.size = Pt(font_size)
            if color:
                run.font.color.rgb = color
            if font_name:
                run.font.name = font_name
        elif part:
            run = p.add_run()
            run.text = part
            run.font.size = Pt(font_size)
            run.font.bold = False
            if color:
                run.font.color.rgb = color
            if font_name:
                run.font.name = font_name
    return txbox


def _create_light_slide(prs: Any, colors: dict[str, Any]) -> Any:
    """Create a slide with light background."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = colors["light"]
    return slide


def _add_slide_heading(
    slide: Any, text: str, colors: dict[str, Any], *, text_color: Any = None
) -> float:
    """Add a title heading with accent bar underneath. Returns y position after heading."""
    from pptx.util import Inches, Pt

    heading_color = text_color or colors["white"]
    font_name = colors.get("font_family")
    _add_text_box(
        slide,
        Inches(0.8),
        Inches(0.5),
        Inches(11),
        Inches(1),
        text,
        font_size=36,
        bold=True,
        color=heading_color,
        font_name=font_name,
    )
    # Accent bar
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(0.8),
        Inches(1.35),
        Inches(2.5),
        Inches(0.06),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = colors["accent"]
    bar.line.fill.background()
    _ = Pt  # suppress unused import warning
    return 2.0


def _add_bullet_list(
    slide: Any,
    left: Any,
    top: Any,
    width: Any,
    items: list[str],
    colors: dict[str, Any],
    *,
    font_size: int = 18,
    spacing: float = 0.6,
    bullet_char: str = "\u2022",
    color: Any = None,
) -> float:
    """Render N bullet items. Returns final y position (inches)."""
    from pptx.util import Inches

    text_color = color or colors["white"]
    font_name = colors.get("font_family")
    y = top if isinstance(top, float) else top / 914400  # convert EMU to inches if needed
    # Normalize to float inches
    if not isinstance(y, float):
        y = float(y)
    left_emu = left if not isinstance(left, float) else Inches(left)
    width_emu = width if not isinstance(width, float) else Inches(width)
    for item in items:
        _add_text_box(
            slide,
            left_emu,
            Inches(y),
            width_emu,
            Inches(spacing),
            f"{bullet_char} {item}",
            font_size=font_size,
            color=text_color,
            font_name=font_name,
        )
        y += spacing
    return y


def _add_table(
    slide: Any,
    left: Any,
    top: Any,
    width: Any,
    headers: list[str],
    rows: list[list[str]],
    colors: dict[str, Any],
    *,
    col_widths: list[float] | None = None,
    font_size: int = 14,
) -> Any:
    """Create a table shape with styled header row and alternating row colors."""
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    row_count = len(rows) + 1  # +1 for header
    col_count = len(headers)
    if col_count == 0:
        return None

    table_shape = slide.shapes.add_table(
        row_count, col_count, left, top, width, Inches(0.4 * row_count)
    )
    table = table_shape.table

    # Set column widths
    if col_widths:
        for i, w in enumerate(col_widths[:col_count]):
            table.columns[i].width = Inches(w)

    font_name = colors.get("font_family")

    # Header row
    for ci, header in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = header
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(font_size)
            paragraph.font.bold = True
            paragraph.font.color.rgb = colors["white"]
            if font_name:
                paragraph.font.name = font_name
        cell.fill.solid()
        cell.fill.fore_color.rgb = colors["accent"]

    # Data rows
    for ri, row in enumerate(rows):
        for ci, value in enumerate(row[:col_count]):
            cell = table.cell(ri + 1, ci)
            cell.text = value
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.color.rgb = colors["dark_text"]
                if font_name:
                    paragraph.font.name = font_name
            # Alternating row colors
            cell.fill.solid()
            if ri % 2 == 0:
                cell.fill.fore_color.rgb = colors["light"]
            else:
                cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)  # type: ignore[no-untyped-call]

    return table_shape


def _add_callout_box(
    slide: Any,
    left: Any,
    top: Any,
    width: Any,
    text: str,
    colors: dict[str, Any],
    *,
    font_size: int = 24,
) -> Any:
    """Dark rounded rect with accent left border and bold text inside."""
    from pptx.util import Inches, Pt

    # Main box
    box_height = Inches(1.2)
    shape = slide.shapes.add_shape(
        5,  # MSO_SHAPE.ROUNDED_RECTANGLE
        left,
        top,
        width,
        box_height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = colors["primary"]
    shape.line.fill.background()

    # Accent left border
    border = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        left,
        top,
        Inches(0.08),
        box_height,
    )
    border.fill.solid()
    border.fill.fore_color.rgb = colors["accent"]
    border.line.fill.background()

    # Text
    font_name = colors.get("font_family")
    _add_text_box(
        slide,
        left + Inches(0.3),
        top + Inches(0.2),
        width - Inches(0.6),
        box_height - Inches(0.4),
        text,
        font_size=font_size,
        bold=True,
        color=colors["white"],
        font_name=font_name,
    )
    _ = Pt  # suppress unused
    return shape


# =============================================================================
# Slide Builders
# =============================================================================


def _build_title_slide(prs: Any, ctx: PitchContext, colors: dict[str, Any]) -> None:
    """Build the title slide."""
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches

    slide = _create_dark_slide(prs, colors)
    name = ctx.spec.company.name
    tagline = ctx.spec.company.tagline or ""

    _add_text_box(
        slide,
        Inches(1),
        Inches(2.5),
        Inches(11),
        Inches(1.5),
        name,
        font_size=48,
        bold=True,
        color=colors["white"],
        alignment=PP_ALIGN.CENTER,
    )
    if tagline:
        _add_text_box(
            slide,
            Inches(1),
            Inches(4.0),
            Inches(11),
            Inches(1),
            tagline,
            font_size=24,
            color=colors["accent"],
            alignment=PP_ALIGN.CENTER,
        )

    stage = ctx.spec.company.stage.value.replace("_", " ").title()
    _add_text_box(
        slide,
        Inches(5),
        Inches(5.5),
        Inches(3),
        Inches(0.5),
        stage,
        font_size=14,
        color=colors["muted"],
        alignment=PP_ALIGN.CENTER,
    )
    notes = ctx.spec.company.speaker_notes if ctx.spec.company else None
    _add_speaker_notes(slide, notes or f"Title slide for {name}. {tagline}")


def _build_problem_slide(prs: Any, ctx: PitchContext, colors: dict[str, Any]) -> None:
    """Build the problem slide."""
    from pptx.util import Inches

    problem = ctx.spec.problem
    if not problem:
        return

    slide = _create_dark_slide(prs, colors)
    y = _add_slide_heading(slide, problem.headline, colors, text_color=colors["highlight"])

    y = _add_bullet_list(
        slide, Inches(1.2), y, Inches(10), problem.points, colors, font_size=20, spacing=0.7
    )

    if problem.market_failure:
        y += 0.3
        font_name = colors.get("font_family")
        _add_text_box(
            slide,
            Inches(0.8),
            Inches(y),
            Inches(11),
            Inches(0.6),
            "Market Failure",
            font_size=24,
            bold=True,
            color=colors["accent"],
            font_name=font_name,
        )
        y += 0.7
        y = _add_bullet_list(
            slide,
            Inches(1.2),
            y,
            Inches(10),
            problem.market_failure,
            colors,
            font_size=18,
            spacing=0.6,
            bullet_char="\u2192",
            color=colors["muted"],
        )

    notes = problem.speaker_notes
    _add_speaker_notes(slide, notes or f"Problem: {problem.headline}")


def _build_solution_slide(prs: Any, ctx: PitchContext, colors: dict[str, Any]) -> None:
    """Build the solution slide."""
    from pptx.util import Inches

    solution = ctx.spec.solution
    if not solution:
        return

    slide = _create_dark_slide(prs, colors)
    y = _add_slide_heading(slide, solution.headline, colors, text_color=colors["success"])

    font_name = colors.get("font_family")
    if solution.how_it_works:
        for i, step in enumerate(solution.how_it_works, 1):
            _add_text_box(
                slide,
                Inches(1.2),
                Inches(y),
                Inches(5),
                Inches(0.6),
                f"{i}. {step}",
                font_size=18,
                color=colors["white"],
                font_name=font_name,
            )
            y += 0.6

    if solution.value_props:
        vy = 2.0
        _add_text_box(
            slide,
            Inches(7),
            Inches(vy - 0.5),
            Inches(5),
            Inches(0.5),
            "Value Propositions",
            font_size=20,
            bold=True,
            color=colors["accent"],
            font_name=font_name,
        )
        for prop in solution.value_props:
            _add_text_box(
                slide,
                Inches(7.2),
                Inches(vy + 0.2),
                Inches(5),
                Inches(0.5),
                f"\u2713 {prop}",
                font_size=16,
                color=colors["white"],
                font_name=font_name,
            )
            vy += 0.5

    notes = solution.speaker_notes
    _add_speaker_notes(slide, notes or f"Solution: {solution.headline}")


def _build_platform_slide(prs: Any, ctx: PitchContext, colors: dict[str, Any]) -> None:
    """Build platform overview slide from DSL data."""
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches

    if not ctx.entities and not ctx.surfaces:
        return

    slide = _create_dark_slide(prs, colors)
    _add_slide_heading(slide, "Platform Overview", colors)

    metrics: list[tuple[str, str]] = []
    if ctx.entities:
        metrics.append((str(len(ctx.entities)), "Data Models"))
    if ctx.surfaces:
        metrics.append((str(len(ctx.surfaces)), "Screens"))
    if ctx.state_machines:
        metrics.append((str(len(ctx.state_machines)), "Workflows"))
    if ctx.story_count:
        metrics.append((str(ctx.story_count), "User Stories"))

    if metrics:
        x_start = 1.0
        box_w = 2.5
        for i, (value, label) in enumerate(metrics[:4]):
            x = x_start + i * 3.0
            _add_text_box(
                slide,
                Inches(x),
                Inches(2.0),
                Inches(box_w),
                Inches(1),
                value,
                font_size=48,
                bold=True,
                color=colors["accent"],
                alignment=PP_ALIGN.CENTER,
            )
            _add_text_box(
                slide,
                Inches(x),
                Inches(3.0),
                Inches(box_w),
                Inches(0.5),
                label,
                font_size=16,
                color=colors["muted"],
                alignment=PP_ALIGN.CENTER,
            )

    if ctx.entities:
        y = 4.2
        _add_text_box(
            slide,
            Inches(0.8),
            Inches(y),
            Inches(5),
            Inches(0.5),
            "Core Entities",
            font_size=18,
            bold=True,
            color=colors["accent"],
        )
        y += 0.5
        entity_text = ", ".join(ctx.entities[:8])
        if len(ctx.entities) > 8:
            entity_text += f" +{len(ctx.entities) - 8} more"
        _add_text_box(
            slide,
            Inches(1.0),
            Inches(y),
            Inches(10),
            Inches(0.5),
            entity_text,
            font_size=14,
            color=colors["white"],
        )

    _add_speaker_notes(
        slide,
        f"Platform built with {len(ctx.entities)} entities, "
        f"{len(ctx.surfaces)} surfaces. Auto-generated from DSL.",
    )


def _build_personas_slide(prs: Any, ctx: PitchContext, colors: dict[str, Any]) -> None:
    """Build personas slide from DSL data."""
    from pptx.util import Inches

    if not ctx.personas:
        return

    slide = _create_dark_slide(prs, colors)
    _add_slide_heading(slide, "User Personas", colors)

    y = 2.0
    for persona in ctx.personas[:6]:
        _add_text_box(
            slide,
            Inches(1.2),
            Inches(y),
            Inches(3),
            Inches(0.5),
            persona["label"],
            font_size=22,
            bold=True,
            color=colors["accent"],
        )
        if persona.get("description"):
            _add_text_box(
                slide,
                Inches(4.5),
                Inches(y),
                Inches(7),
                Inches(0.5),
                persona["description"],
                font_size=16,
                color=colors["muted"],
            )
        y += 0.8

    _add_speaker_notes(slide, f"{len(ctx.personas)} user personas defined in DSL.")


def _build_market_slide(prs: Any, ctx: PitchContext, colors: dict[str, Any]) -> None:
    """Build market sizing slide."""
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches

    market = ctx.spec.market
    if not market:
        return

    slide = _create_dark_slide(prs, colors)
    _add_slide_heading(slide, "Market Opportunity", colors)

    currency = ctx.spec.company.currency
    sizes = []
    if market.tam:
        sizes.append(("TAM", market.tam))
    if market.sam:
        sizes.append(("SAM", market.sam))
    if market.som:
        sizes.append(("SOM", market.som))

    x_start = 1.0
    for i, (label, size) in enumerate(sizes):
        x = x_start + i * 4.0
        _add_text_box(
            slide,
            Inches(x),
            Inches(2.0),
            Inches(3.5),
            Inches(0.5),
            label,
            font_size=20,
            bold=True,
            color=colors["muted"],
            alignment=PP_ALIGN.CENTER,
        )
        _add_text_box(
            slide,
            Inches(x),
            Inches(2.5),
            Inches(3.5),
            Inches(1),
            _fmt_currency(size.value, currency),
            font_size=42,
            bold=True,
            color=colors["accent"],
            alignment=PP_ALIGN.CENTER,
        )
        _add_text_box(
            slide,
            Inches(x),
            Inches(3.5),
            Inches(3.5),
            Inches(0.5),
            size.label,
            font_size=14,
            color=colors["white"],
            alignment=PP_ALIGN.CENTER,
        )

    if market.drivers:
        y = 5.0
        _add_text_box(
            slide,
            Inches(0.8),
            Inches(y),
            Inches(11),
            Inches(0.5),
            "Market Drivers",
            font_size=20,
            bold=True,
            color=colors["accent"],
        )
        y += 0.6
        for driver in market.drivers:
            _add_text_box(
                slide,
                Inches(1.2),
                Inches(y),
                Inches(10),
                Inches(0.4),
                f"\u25b8 {driver}",
                font_size=16,
                color=colors["white"],
            )
            y += 0.5

    # Embed market chart if available
    market_chart = ctx.chart_paths.get("market")
    if market_chart and Path(market_chart).exists():
        try:
            slide.shapes.add_picture(
                str(market_chart), Inches(7), Inches(4.5), Inches(5), Inches(2.5)
            )
        except Exception as e:
            logger.debug(f"Failed to embed market chart: {e}")

    notes = market.speaker_notes
    _add_speaker_notes(slide, notes or "Market sizing: TAM/SAM/SOM breakdown.")


def _build_business_model_slide(prs: Any, ctx: PitchContext, colors: dict[str, Any]) -> None:
    """Build business model / pricing slide."""
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches

    bm = ctx.spec.business_model
    if not bm or not bm.tiers:
        return

    slide = _create_dark_slide(prs, colors)
    _add_slide_heading(slide, "Business Model", colors)

    currency = ctx.spec.company.currency
    tier_count = len(bm.tiers)
    box_w = min(3.5, 11.0 / max(tier_count, 1))

    for i, tier in enumerate(bm.tiers):
        x = 0.8 + i * (box_w + 0.3)
        y = 2.0

        name_color = colors["highlight"] if tier.highlighted else colors["white"]
        _add_text_box(
            slide,
            Inches(x),
            Inches(y),
            Inches(box_w),
            Inches(0.5),
            tier.name,
            font_size=22,
            bold=True,
            color=name_color,
            alignment=PP_ALIGN.CENTER,
        )

        price_str = _fmt_currency(tier.price, currency) if tier.price else "Free"
        _add_text_box(
            slide,
            Inches(x),
            Inches(y + 0.6),
            Inches(box_w),
            Inches(0.8),
            f"{price_str}/{tier.period}",
            font_size=28,
            bold=True,
            color=colors["accent"],
            alignment=PP_ALIGN.CENTER,
        )

        if tier.features:
            _add_text_box(
                slide,
                Inches(x),
                Inches(y + 1.5),
                Inches(box_w),
                Inches(2),
                tier.features,
                font_size=14,
                color=colors["muted"],
                alignment=PP_ALIGN.CENTER,
            )

    notes = bm.speaker_notes
    _add_speaker_notes(slide, notes or f"Business model with {tier_count} pricing tiers.")


def _build_financials_slide(prs: Any, ctx: PitchContext, colors: dict[str, Any]) -> None:
    """Build financials slide with projections."""
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches

    fin = ctx.spec.financials
    if not fin or not fin.projections:
        return

    slide = _create_dark_slide(prs, colors)
    _add_slide_heading(slide, "Financial Projections", colors)

    currency = ctx.spec.company.currency
    col_count = len(fin.projections)
    col_w = min(3.0, 11.0 / max(col_count, 1))

    for i, proj in enumerate(fin.projections):
        x = 0.8 + i * (col_w + 0.3)
        _add_text_box(
            slide,
            Inches(x),
            Inches(2.0),
            Inches(col_w),
            Inches(0.5),
            str(proj.year),
            font_size=20,
            bold=True,
            color=colors["muted"],
            alignment=PP_ALIGN.CENTER,
        )
        _add_text_box(
            slide,
            Inches(x),
            Inches(2.6),
            Inches(col_w),
            Inches(0.8),
            _fmt_currency(proj.revenue, currency),
            font_size=32,
            bold=True,
            color=colors["success"],
            alignment=PP_ALIGN.CENTER,
        )
        _add_text_box(
            slide,
            Inches(x),
            Inches(3.4),
            Inches(col_w),
            Inches(0.4),
            "Revenue",
            font_size=12,
            color=colors["muted"],
            alignment=PP_ALIGN.CENTER,
        )
        _add_text_box(
            slide,
            Inches(x),
            Inches(4.0),
            Inches(col_w),
            Inches(0.5),
            f"{proj.customers:,} customers",
            font_size=16,
            color=colors["white"],
            alignment=PP_ALIGN.CENTER,
        )

    if fin.use_of_funds:
        y = 5.2
        _add_text_box(
            slide,
            Inches(0.8),
            Inches(y),
            Inches(11),
            Inches(0.5),
            "Use of Funds",
            font_size=20,
            bold=True,
            color=colors["accent"],
        )
        y += 0.5
        parts = [f"{f.category} ({f.percent}%)" for f in fin.use_of_funds]
        _add_text_box(
            slide,
            Inches(1.0),
            Inches(y),
            Inches(11),
            Inches(0.5),
            "  |  ".join(parts),
            font_size=14,
            color=colors["white"],
        )

    # Embed revenue chart if available
    revenue_chart = ctx.chart_paths.get("revenue")
    if revenue_chart and Path(revenue_chart).exists():
        try:
            slide.shapes.add_picture(
                str(revenue_chart), Inches(0.8), Inches(4.5), Inches(5), Inches(2.5)
            )
        except Exception as e:
            logger.debug(f"Failed to embed revenue chart: {e}")

    notes = fin.speaker_notes
    _add_speaker_notes(slide, notes or "Financial projections and use of funds breakdown.")


def _build_team_slide(prs: Any, ctx: PitchContext, colors: dict[str, Any]) -> None:
    """Build team slide with light background."""
    from pptx.util import Inches

    team = ctx.spec.team
    if not team:
        return

    slide = _create_light_slide(prs, colors)
    y = _add_slide_heading(slide, "Team", colors, text_color=colors["dark_text"])

    font_name = colors.get("font_family")
    for member in team.founders:
        _add_text_box(
            slide,
            Inches(1.2),
            Inches(y),
            Inches(3),
            Inches(0.5),
            member.name,
            font_size=22,
            bold=True,
            color=colors["accent"],
            font_name=font_name,
        )
        _add_text_box(
            slide,
            Inches(4.5),
            Inches(y),
            Inches(2),
            Inches(0.5),
            member.role,
            font_size=18,
            color=colors["highlight"],
            font_name=font_name,
        )
        if member.bio:
            _add_text_box(
                slide,
                Inches(1.2),
                Inches(y + 0.4),
                Inches(10),
                Inches(0.4),
                member.bio,
                font_size=14,
                color=colors["muted"],
                font_name=font_name,
            )
            y += 1.0
        else:
            y += 0.7

    if team.advisors:
        y += 0.3
        _add_text_box(
            slide,
            Inches(0.8),
            Inches(y),
            Inches(11),
            Inches(0.5),
            "Advisors",
            font_size=20,
            bold=True,
            color=colors["accent"],
            font_name=font_name,
        )
        y += 0.5
        for advisor in team.advisors:
            _add_text_box(
                slide,
                Inches(1.2),
                Inches(y),
                Inches(5),
                Inches(0.4),
                f"{advisor.name} \u2014 {advisor.role}",
                font_size=16,
                color=colors["dark_text"],
                font_name=font_name,
            )
            y += 0.5

    if team.key_hires:
        y += 0.3
        _add_text_box(
            slide,
            Inches(0.8),
            Inches(y),
            Inches(11),
            Inches(0.5),
            "Key Hires",
            font_size=20,
            bold=True,
            color=colors["accent"],
            font_name=font_name,
        )
        y += 0.5
        for hire in team.key_hires:
            timing = f" ({hire.timing})" if hire.timing else ""
            _add_text_box(
                slide,
                Inches(1.2),
                Inches(y),
                Inches(10),
                Inches(0.4),
                f"{hire.role}{timing}",
                font_size=16,
                color=colors["dark_text"],
                font_name=font_name,
            )
            y += 0.5

    notes = team.speaker_notes
    _add_speaker_notes(
        slide,
        notes or f"Team: {len(team.founders)} founders, {len(team.advisors)} advisors.",
    )


def _build_competition_slide(prs: Any, ctx: PitchContext, colors: dict[str, Any]) -> None:
    """Build competition slide with light background and table."""
    from pptx.util import Inches

    if not ctx.spec.competitors:
        return

    slide = _create_light_slide(prs, colors)
    _add_slide_heading(slide, "Competitive Landscape", colors, text_color=colors["dark_text"])

    headers = ["Competitor", "Strength", "Weakness"]
    rows = [
        [comp.name, comp.strength or "", comp.weakness or ""] for comp in ctx.spec.competitors[:6]
    ]
    _add_table(
        slide,
        Inches(0.8),
        Inches(2.2),
        Inches(11.5),
        headers,
        rows,
        colors,
        col_widths=[3.5, 4.0, 4.0],
        font_size=14,
    )

    _add_speaker_notes(slide, f"{len(ctx.spec.competitors)} competitors analyzed.")


def _build_milestones_slide(prs: Any, ctx: PitchContext, colors: dict[str, Any]) -> None:
    """Build milestones / roadmap slide."""
    from pptx.util import Inches

    ms = ctx.spec.milestones
    if not ms:
        return

    slide = _create_dark_slide(prs, colors)
    y = _add_slide_heading(slide, "Milestones & Roadmap", colors)

    font_name = colors.get("font_family")

    if ms.completed:
        _add_text_box(
            slide,
            Inches(0.8),
            Inches(y),
            Inches(5),
            Inches(0.5),
            "Completed \u2713",
            font_size=20,
            bold=True,
            color=colors["success"],
            font_name=font_name,
        )
        y += 0.5
        y = _add_bullet_list(
            slide,
            Inches(1.2),
            y,
            Inches(10),
            ms.completed,
            colors,
            font_size=16,
            spacing=0.5,
            bullet_char="\u2713",
            color=colors["muted"],
        )

    if ms.next_12_months:
        y += 0.3
        _add_text_box(
            slide,
            Inches(0.8),
            Inches(y),
            Inches(5),
            Inches(0.5),
            "Next 12 Months",
            font_size=20,
            bold=True,
            color=colors["accent"],
            font_name=font_name,
        )
        y += 0.5
        y = _add_bullet_list(
            slide,
            Inches(1.2),
            y,
            Inches(10),
            ms.next_12_months,
            colors,
            font_size=16,
            spacing=0.5,
            bullet_char="\u2192",
        )

    if ms.long_term:
        y += 0.3
        _add_text_box(
            slide,
            Inches(0.8),
            Inches(y),
            Inches(5),
            Inches(0.5),
            "Long Term Vision",
            font_size=20,
            bold=True,
            color=colors["highlight"],
            font_name=font_name,
        )
        y += 0.5
        y = _add_bullet_list(
            slide,
            Inches(1.2),
            y,
            Inches(10),
            ms.long_term,
            colors,
            font_size=16,
            spacing=0.5,
            bullet_char="\u25c6",
            color=colors["muted"],
        )

    notes = ms.speaker_notes
    _add_speaker_notes(slide, notes or "Milestones and roadmap.")


def _build_ask_slide(prs: Any, ctx: PitchContext, colors: dict[str, Any]) -> None:
    """Build the funding ask slide."""
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches

    if not ctx.spec.company.funding_ask:
        return

    slide = _create_dark_slide(prs, colors)
    _add_slide_heading(slide, "The Ask", colors)

    currency = ctx.spec.company.currency
    ask_str = _fmt_currency(ctx.spec.company.funding_ask, currency)
    stage = ctx.spec.company.stage.value.replace("_", " ").title()

    _add_text_box(
        slide,
        Inches(1),
        Inches(2.5),
        Inches(11),
        Inches(1.5),
        ask_str,
        font_size=64,
        bold=True,
        color=colors["highlight"],
        alignment=PP_ALIGN.CENTER,
    )
    _add_text_box(
        slide,
        Inches(1),
        Inches(4.0),
        Inches(11),
        Inches(0.5),
        f"{stage} Round",
        font_size=24,
        color=colors["accent"],
        alignment=PP_ALIGN.CENTER,
    )

    if ctx.spec.company.runway_months:
        _add_text_box(
            slide,
            Inches(1),
            Inches(5.0),
            Inches(11),
            Inches(0.5),
            f"{ctx.spec.company.runway_months} months runway",
            font_size=18,
            color=colors["muted"],
            alignment=PP_ALIGN.CENTER,
        )

    if ctx.spec.financials and ctx.spec.financials.use_of_funds:
        y = 5.8
        parts = [f"{f.category} ({f.percent}%)" for f in ctx.spec.financials.use_of_funds]
        _add_text_box(
            slide,
            Inches(1),
            Inches(y),
            Inches(11),
            Inches(0.5),
            "  |  ".join(parts),
            font_size=14,
            color=colors["white"],
            alignment=PP_ALIGN.CENTER,
        )

    # Embed funds chart if available
    funds_chart = ctx.chart_paths.get("funds")
    if funds_chart and Path(funds_chart).exists():
        try:
            slide.shapes.add_picture(
                str(funds_chart), Inches(8), Inches(4.5), Inches(4), Inches(2.5)
            )
        except Exception as e:
            logger.debug(f"Failed to embed funds chart: {e}")

    _add_speaker_notes(slide, f"Raising {ask_str} at {stage} stage.")


def _build_closing_slide(prs: Any, ctx: PitchContext, colors: dict[str, Any]) -> None:
    """Build the closing / thank you slide."""
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches

    slide = _create_dark_slide(prs, colors)
    _add_text_box(
        slide,
        Inches(1),
        Inches(2.5),
        Inches(11),
        Inches(1.5),
        "Thank You",
        font_size=48,
        bold=True,
        color=colors["white"],
        alignment=PP_ALIGN.CENTER,
    )
    _add_text_box(
        slide,
        Inches(1),
        Inches(4.2),
        Inches(11),
        Inches(0.8),
        ctx.spec.company.name,
        font_size=28,
        color=colors["accent"],
        alignment=PP_ALIGN.CENTER,
    )
    if ctx.spec.company.tagline:
        _add_text_box(
            slide,
            Inches(1),
            Inches(5.0),
            Inches(11),
            Inches(0.5),
            ctx.spec.company.tagline,
            font_size=18,
            color=colors["muted"],
            alignment=PP_ALIGN.CENTER,
        )

    notes = ctx.spec.company.speaker_notes if ctx.spec.company else None
    _add_speaker_notes(slide, notes or "Closing slide. Open for questions.")


def _build_extra_slide(
    prs: Any, ctx: PitchContext, colors: dict[str, Any], extra: ExtraSlide
) -> None:
    """Build a user-defined extra slide based on its layout."""
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches

    use_light = extra.theme == "light"
    slide = _create_light_slide(prs, colors) if use_light else _create_dark_slide(prs, colors)
    heading_color = colors["dark_text"] if use_light else colors["white"]
    text_color = colors["dark_text"] if use_light else colors["white"]
    y = _add_slide_heading(slide, extra.title, colors, text_color=heading_color)

    font_name = colors.get("font_family")

    if extra.layout == ExtraSlideLayout.BULLETS:
        _add_bullet_list(
            slide,
            Inches(1.2),
            y,
            Inches(10),
            extra.items,
            colors,
            font_size=20,
            spacing=0.7,
            color=text_color,
        )

    elif extra.layout == ExtraSlideLayout.STATS:
        parsed: list[tuple[str, str]] = []
        for item in extra.items:
            if "|" in item:
                val, lbl = item.split("|", 1)
                parsed.append((val.strip(), lbl.strip()))
            else:
                parsed.append((item, ""))
        _add_columns(
            slide,
            Inches(2.0),
            parsed,
            value_color=colors["accent"],
            label_color=colors["muted"],
            alignment=PP_ALIGN.CENTER,
        )

    elif extra.layout == ExtraSlideLayout.CARDS:
        # Multi-column card grid
        items = extra.items
        col_count = min(3, len(items)) if items else 1
        margin = 1.0
        gap = 0.3
        card_w = (11.0 - gap * (col_count - 1)) / col_count
        card_h = 0.8
        row_gap = 0.3
        for idx, item in enumerate(items):
            col = idx % col_count
            row = idx // col_count
            x = margin + col * (card_w + gap)
            cy = y + row * (card_h + row_gap)
            _add_card(
                slide,
                Inches(x),
                Inches(cy),
                Inches(card_w),
                Inches(card_h),
                fill_color=colors.get("dark_text"),
                border_color=colors["accent"],
            )
            _add_rich_text_box(
                slide,
                Inches(x + 0.3),
                Inches(cy + 0.1),
                Inches(card_w - 0.6),
                Inches(0.6),
                item,
                font_size=18,
                color=colors["white"] if not use_light else colors["white"],
                font_name=font_name,
            )

    elif extra.layout == ExtraSlideLayout.TABLE:
        # Parse items as pipe-separated: first item = headers, rest = rows
        if extra.items:
            headers = [h.strip() for h in extra.items[0].split("|")]
            rows = [[c.strip() for c in item.split("|")] for item in extra.items[1:]]
            _add_table(
                slide,
                Inches(0.8),
                Inches(y),
                Inches(11.5),
                headers,
                rows,
                colors,
                font_size=14,
            )

    elif extra.layout == ExtraSlideLayout.CALLOUT:
        # First item as callout box, rest as supporting bullets
        if extra.items:
            _add_callout_box(
                slide,
                Inches(0.8),
                Inches(y),
                Inches(11.5),
                extra.items[0],
                colors,
                font_size=24,
            )
            if len(extra.items) > 1:
                _add_bullet_list(
                    slide,
                    Inches(1.2),
                    y + 1.5,
                    Inches(10),
                    extra.items[1:],
                    colors,
                    font_size=18,
                    spacing=0.6,
                    color=text_color,
                )

    elif extra.layout == ExtraSlideLayout.IMAGE:
        if extra.image_path:
            try:
                slide.shapes.add_picture(
                    extra.image_path,
                    Inches(1.5),
                    Inches(2.0),
                    Inches(10),
                    Inches(5),
                )
            except Exception as e:
                logger.warning(f"Failed to embed image {extra.image_path}: {e}")

    notes = extra.speaker_notes
    _add_speaker_notes(slide, notes or f"Extra slide: {extra.title}")


# =============================================================================
# Slide Catalog
# =============================================================================


def _has_problem(ctx: PitchContext) -> bool:
    return ctx.spec.problem is not None


def _has_solution(ctx: PitchContext) -> bool:
    return ctx.spec.solution is not None


def _has_dsl(ctx: PitchContext) -> bool:
    return bool(ctx.entities or ctx.surfaces)


def _has_personas(ctx: PitchContext) -> bool:
    return bool(ctx.personas)


def _has_market(ctx: PitchContext) -> bool:
    return ctx.spec.market is not None


def _has_tiers(ctx: PitchContext) -> bool:
    return ctx.spec.business_model is not None and bool(ctx.spec.business_model.tiers)


def _has_projections(ctx: PitchContext) -> bool:
    return ctx.spec.financials is not None and bool(ctx.spec.financials.projections)


def _has_team(ctx: PitchContext) -> bool:
    return ctx.spec.team is not None


def _has_competitors(ctx: PitchContext) -> bool:
    return bool(ctx.spec.competitors)


def _has_milestones(ctx: PitchContext) -> bool:
    return ctx.spec.milestones is not None


def _has_funding_ask(ctx: PitchContext) -> bool:
    return ctx.spec.company.funding_ask is not None


def _always(_ctx: PitchContext) -> bool:
    return True


# Ordered slide catalog: (name, builder, condition)
SlideCatalogEntry = tuple[str, Callable[..., None], Callable[[PitchContext], bool]]

SLIDE_CATALOG: list[SlideCatalogEntry] = [
    ("title", _build_title_slide, _always),
    ("problem", _build_problem_slide, _has_problem),
    ("solution", _build_solution_slide, _has_solution),
    ("platform", _build_platform_slide, _has_dsl),
    ("personas", _build_personas_slide, _has_personas),
    ("market", _build_market_slide, _has_market),
    ("business_model", _build_business_model_slide, _has_tiers),
    ("financials", _build_financials_slide, _has_projections),
    ("team", _build_team_slide, _has_team),
    ("competition", _build_competition_slide, _has_competitors),
    ("milestones", _build_milestones_slide, _has_milestones),
    ("ask", _build_ask_slide, _has_funding_ask),
    ("closing", _build_closing_slide, _always),
]


# =============================================================================
# Main Generator
# =============================================================================


def generate_pptx(ctx: PitchContext, output_path: Path) -> GeneratorResult:
    """Generate PPTX pitch deck from PitchContext.

    Args:
        ctx: PitchContext with merged PitchSpec + DSL data.
        output_path: Path to write the .pptx file.

    Returns:
        GeneratorResult with success status and output path.
    """
    if not _check_pptx_available():
        return GeneratorResult(
            success=False,
            error="python-pptx is not installed. Install with: pip install 'dazzle[pitch]'",
        )

    try:
        from pptx import Presentation as PptxPresentation
        from pptx.util import Inches

        prs = PptxPresentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        colors = _resolve_colors(ctx.spec.brand)

        # Generate charts and populate ctx.chart_paths
        try:
            from dazzle.pitch.generators.charts import (
                generate_funds_chart,
                generate_market_chart,
                generate_revenue_chart,
            )

            chart_dir = output_path.parent / "charts"
            chart_dir.mkdir(parents=True, exist_ok=True)
            hex_colors = {
                "primary": ctx.spec.brand.primary,
                "accent": ctx.spec.brand.accent,
                "highlight": ctx.spec.brand.highlight,
                "success": ctx.spec.brand.success,
                "light": ctx.spec.brand.light,
            }
            for name, gen_fn in [
                ("revenue", generate_revenue_chart),
                ("market", generate_market_chart),
                ("funds", generate_funds_chart),
            ]:
                path = gen_fn(ctx, chart_dir, hex_colors)
                if path:
                    ctx.chart_paths[name] = path
        except Exception as e:
            logger.debug(f"Chart generation skipped: {e}")

        # Build catalog and extra slide maps
        catalog_map: dict[str, tuple[Callable[..., None], Callable[[PitchContext], bool]]] = {
            entry_name: (builder, condition) for entry_name, builder, condition in SLIDE_CATALOG
        }
        extra_map: dict[str, ExtraSlide] = {
            es.title.lower().replace(" ", "_"): es for es in ctx.spec.extra_slides
        }

        # Determine slide order
        if ctx.spec.slide_order:
            ordered_names = list(ctx.spec.slide_order)
        else:
            ordered_names = [entry_name for entry_name, _, _ in SLIDE_CATALOG]
            # Insert extras before closing
            if extra_map:
                closing_idx = (
                    ordered_names.index("closing")
                    if "closing" in ordered_names
                    else len(ordered_names)
                )
                for extra_name in extra_map:
                    ordered_names.insert(closing_idx, extra_name)
                    closing_idx += 1

        slide_count = 0
        for slide_name in ordered_names:
            if slide_name in catalog_map:
                builder, condition = catalog_map[slide_name]
                if condition(ctx):
                    builder(prs, ctx, colors)
                    slide_count += 1
                    logger.debug(f"Built slide: {slide_name}")
            elif slide_name in extra_map:
                _build_extra_slide(prs, ctx, colors, extra_map[slide_name])
                slide_count += 1
                logger.debug(f"Built extra slide: {slide_name}")

        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))

        return GeneratorResult(
            success=True,
            output_path=output_path,
            files_created=[str(output_path)],
            slide_count=slide_count,
        )

    except Exception as e:
        logger.exception("Error generating PPTX")
        return GeneratorResult(success=False, error=str(e))
