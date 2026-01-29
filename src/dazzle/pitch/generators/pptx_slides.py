"""
PPTX slide builder functions.

Each function builds one slide type from PitchContext data.
"""

from __future__ import annotations

import logging
from pathlib import Path
from typing import Any

from dazzle.pitch.extractor import PitchContext
from dazzle.pitch.generators.pptx_primitives import (
    CONTENT_BOTTOM,
    ContentRegion,
    _add_bullet_list,
    _add_callout_box,
    _add_card,
    _add_columns,
    _add_divider,
    _add_rich_text_box,
    _add_slide_heading,
    _add_speaker_notes,
    _add_stat_box,
    _add_table,
    _add_text_box,
    _create_dark_slide,
    _create_light_slide,
    _estimate_text_height,
    _fmt_currency,
)
from dazzle.pitch.ir import ExtraSlide, ExtraSlideLayout

logger = logging.getLogger(__name__)

# Re-export stat_box and columns for type-checker satisfaction
__all__ = [
    "_build_title_slide",
    "_build_problem_slide",
    "_build_solution_slide",
    "_build_platform_slide",
    "_build_personas_slide",
    "_build_market_slide",
    "_build_business_model_slide",
    "_build_financials_slide",
    "_build_team_slide",
    "_build_competition_slide",
    "_build_milestones_slide",
    "_build_ask_slide",
    "_build_closing_slide",
    "_build_extra_slide",
]

# Suppress unused-import warnings for primitives used only by specific builders
_ = (_add_stat_box, _add_columns, _add_rich_text_box, _add_callout_box, _add_table, _add_divider)


def _build_title_slide(prs: Any, ctx: PitchContext, colors: dict[str, Any]) -> None:
    """Build the title slide."""
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches

    slide = _create_dark_slide(prs, colors)
    name = ctx.spec.company.name
    tagline = ctx.spec.company.tagline or ""

    _add_text_box(
        slide,
        Inches(1),
        Inches(2.5),
        Inches(11),
        Inches(1.5),
        name,
        font_size=48,
        bold=True,
        color=colors["white"],
        alignment=PP_ALIGN.CENTER,
    )
    if tagline:
        _add_text_box(
            slide,
            Inches(1),
            Inches(4.0),
            Inches(11),
            Inches(1),
            tagline,
            font_size=24,
            color=colors["accent"],
            alignment=PP_ALIGN.CENTER,
        )

    stage = ctx.spec.company.stage.value.replace("_", " ").title()
    _add_text_box(
        slide,
        Inches(5),
        Inches(5.5),
        Inches(3),
        Inches(0.5),
        stage,
        font_size=14,
        color=colors["muted"],
        alignment=PP_ALIGN.CENTER,
    )

    # Logo embedding
    if ctx.spec.company.logo_path:
        logo = Path(ctx.spec.company.logo_path)
        if logo.exists():
            slide.shapes.add_picture(str(logo), Inches(5.5), Inches(0.5), Inches(2), Inches(1))

    notes = ctx.spec.company.speaker_notes if ctx.spec.company else None
    _add_speaker_notes(slide, notes or f"Title slide for {name}. {tagline}")


def _build_problem_slide(prs: Any, ctx: PitchContext, colors: dict[str, Any]) -> None:
    """Build the problem slide."""
    from pptx.util import Inches

    problem = ctx.spec.problem
    if not problem:
        return

    slide = _create_dark_slide(prs, colors)
    y = _add_slide_heading(slide, problem.headline, colors, text_color=colors["highlight"])

    y = _add_bullet_list(
        slide, Inches(1.2), y, Inches(10), problem.points, colors, font_size=20, spacing=0.7
    ).final_y

    if problem.market_failure:
        y += 0.3
        font_name = colors.get("font_family")
        _add_text_box(
            slide,
            Inches(0.8),
            Inches(y),
            Inches(11),
            Inches(0.6),
            "Market Failure",
            font_size=24,
            bold=True,
            color=colors["accent"],
            font_name=font_name,
        )
        y += 0.7
        y = _add_bullet_list(
            slide,
            Inches(1.2),
            y,
            Inches(10),
            problem.market_failure,
            colors,
            font_size=18,
            spacing=0.6,
            bullet_char="\u2192",
            color=colors["muted"],
        ).final_y

    notes = problem.speaker_notes
    _add_speaker_notes(slide, notes or f"Problem: {problem.headline}")


def _build_solution_slide(prs: Any, ctx: PitchContext, colors: dict[str, Any]) -> None:
    """Build the solution slide."""
    from pptx.util import Inches

    solution = ctx.spec.solution
    if not solution:
        return

    slide = _create_dark_slide(prs, colors)
    y = _add_slide_heading(slide, solution.headline, colors, text_color=colors["success"])

    font_name = colors.get("font_family")
    if solution.how_it_works:
        left_region = ContentRegion(left=1.2, top=y, width=5.0, bottom=CONTENT_BOTTOM)
        for i, step in enumerate(solution.how_it_works, 1):
            step_text = f"{i}. {step}"
            est_h = max(0.6, _estimate_text_height(step_text, 5.0, 18))
            if not left_region.fits(est_h):
                remaining = len(solution.how_it_works) - i + 1
                logger.warning(f"Solution slide: truncated {remaining} how_it_works steps")
                break
            _add_text_box(
                slide,
                Inches(left_region.left),
                Inches(left_region.top),
                Inches(5),
                Inches(est_h),
                step_text,
                font_size=18,
                color=colors["white"],
                font_name=font_name,
            )
            left_region = left_region.advance(est_h)

    if solution.value_props:
        right_region = ContentRegion(left=7.0, top=2.0, width=5.0, bottom=CONTENT_BOTTOM)
        _add_text_box(
            slide,
            Inches(7),
            Inches(right_region.top - 0.5),
            Inches(5),
            Inches(0.5),
            "Value Propositions",
            font_size=20,
            bold=True,
            color=colors["accent"],
            font_name=font_name,
        )
        right_region = right_region.advance(0.2)
        for prop in solution.value_props:
            prop_text = f"\u2713 {prop}"
            est_h = max(0.5, _estimate_text_height(prop_text, 5.0, 16))
            if not right_region.fits(est_h):
                logger.warning("Solution slide: truncated value_props")
                break
            _add_text_box(
                slide,
                Inches(7.2),
                Inches(right_region.top),
                Inches(5),
                Inches(est_h),
                prop_text,
                font_size=16,
                color=colors["white"],
                font_name=font_name,
            )
            right_region = right_region.advance(est_h)

    notes = solution.speaker_notes
    _add_speaker_notes(slide, notes or f"Solution: {solution.headline}")


def _build_platform_slide(prs: Any, ctx: PitchContext, colors: dict[str, Any]) -> None:
    """Build platform overview slide from DSL data."""
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches

    if not ctx.entities and not ctx.surfaces:
        return

    slide = _create_dark_slide(prs, colors)
    _add_slide_heading(slide, "Platform Overview", colors)

    metrics: list[tuple[str, str]] = []
    if ctx.entities:
        metrics.append((str(len(ctx.entities)), "Data Models"))
    if ctx.surfaces:
        metrics.append((str(len(ctx.surfaces)), "Screens"))
    if ctx.state_machines:
        metrics.append((str(len(ctx.state_machines)), "Workflows"))
    if ctx.story_count:
        metrics.append((str(ctx.story_count), "User Stories"))

    if metrics:
        x_start = 1.0
        box_w = 2.5
        for i, (value, label) in enumerate(metrics[:4]):
            x = x_start + i * 3.0
            # Card background behind each stat
            _add_card(
                slide,
                Inches(x - 0.2),
                Inches(1.8),
                Inches(box_w + 0.4),
                Inches(1.8),
                fill_color=colors.get("dark_text"),
                border_color=colors["accent"],
            )
            _add_text_box(
                slide,
                Inches(x),
                Inches(2.0),
                Inches(box_w),
                Inches(1),
                value,
                font_size=48,
                bold=True,
                color=colors["accent"],
                alignment=PP_ALIGN.CENTER,
            )
            _add_text_box(
                slide,
                Inches(x),
                Inches(3.0),
                Inches(box_w),
                Inches(0.5),
                label,
                font_size=16,
                color=colors["muted"],
                alignment=PP_ALIGN.CENTER,
            )

    if ctx.entities:
        y = 4.2
        _add_text_box(
            slide,
            Inches(0.8),
            Inches(y),
            Inches(5),
            Inches(0.5),
            "Core Entities",
            font_size=18,
            bold=True,
            color=colors["accent"],
        )
        y += 0.5
        entity_text = ", ".join(ctx.entities[:8])
        if len(ctx.entities) > 8:
            entity_text += f" +{len(ctx.entities) - 8} more"
        _add_text_box(
            slide,
            Inches(1.0),
            Inches(y),
            Inches(10),
            Inches(0.5),
            entity_text,
            font_size=14,
            color=colors["white"],
        )

    _add_speaker_notes(
        slide,
        f"Platform built with {len(ctx.entities)} entities, "
        f"{len(ctx.surfaces)} surfaces. Auto-generated from DSL.",
    )


def _build_personas_slide(prs: Any, ctx: PitchContext, colors: dict[str, Any]) -> None:
    """Build personas slide from DSL data."""
    from pptx.util import Inches

    if not ctx.personas:
        return

    slide = _create_dark_slide(prs, colors)
    _add_slide_heading(slide, "User Personas", colors)

    region = ContentRegion(left=1.2, top=2.0, width=10.0, bottom=CONTENT_BOTTOM)
    for persona in ctx.personas[:6]:
        if not region.fits(0.8):
            logger.warning("Personas slide: truncated due to overflow")
            break
        _add_text_box(
            slide,
            Inches(1.2),
            Inches(region.top),
            Inches(3),
            Inches(0.5),
            persona["label"],
            font_size=22,
            bold=True,
            color=colors["accent"],
        )
        if persona.get("description"):
            _add_text_box(
                slide,
                Inches(4.5),
                Inches(region.top),
                Inches(7),
                Inches(0.5),
                persona["description"],
                font_size=16,
                color=colors["muted"],
            )
        region = region.advance(0.8)

    _add_speaker_notes(slide, f"{len(ctx.personas)} user personas defined in DSL.")


def _build_market_slide(prs: Any, ctx: PitchContext, colors: dict[str, Any]) -> None:
    """Build market sizing slide."""
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches

    market = ctx.spec.market
    if not market:
        return

    slide = _create_dark_slide(prs, colors)
    _add_slide_heading(slide, "Market Opportunity", colors)

    currency = ctx.spec.company.currency
    sizes = []
    if market.tam:
        sizes.append(("TAM", market.tam))
    if market.sam:
        sizes.append(("SAM", market.sam))
    if market.som:
        sizes.append(("SOM", market.som))

    x_start = 1.0
    for i, (label, size) in enumerate(sizes):
        x = x_start + i * 4.0
        _add_text_box(
            slide,
            Inches(x),
            Inches(2.0),
            Inches(3.5),
            Inches(0.5),
            label,
            font_size=20,
            bold=True,
            color=colors["muted"],
            alignment=PP_ALIGN.CENTER,
        )
        _add_text_box(
            slide,
            Inches(x),
            Inches(2.5),
            Inches(3.5),
            Inches(1),
            _fmt_currency(size.value, currency),
            font_size=42,
            bold=True,
            color=colors["accent"],
            alignment=PP_ALIGN.CENTER,
        )
        _add_text_box(
            slide,
            Inches(x),
            Inches(3.5),
            Inches(3.5),
            Inches(0.5),
            size.label,
            font_size=14,
            color=colors["white"],
            alignment=PP_ALIGN.CENTER,
        )

    if market.drivers:
        driver_region = ContentRegion(left=0.8, top=5.0, width=11.0, bottom=CONTENT_BOTTOM)
        _add_text_box(
            slide,
            Inches(0.8),
            Inches(driver_region.top),
            Inches(11),
            Inches(0.5),
            "Market Drivers",
            font_size=20,
            bold=True,
            color=colors["accent"],
        )
        driver_region = driver_region.advance(0.6)
        for driver in market.drivers:
            if not driver_region.fits(0.5):
                logger.warning("Market slide: truncated drivers due to overflow")
                break
            _add_text_box(
                slide,
                Inches(1.2),
                Inches(driver_region.top),
                Inches(10),
                Inches(0.4),
                f"\u25b8 {driver}",
                font_size=16,
                color=colors["white"],
            )
            driver_region = driver_region.advance(0.5)

    # Embed market chart if available
    market_chart = ctx.chart_paths.get("market")
    if market_chart and Path(market_chart).exists():
        try:
            slide.shapes.add_picture(
                str(market_chart), Inches(7), Inches(4.5), Inches(5), Inches(2.5)
            )
        except Exception as e:
            logger.debug(f"Failed to embed market chart: {e}")

    notes = market.speaker_notes
    _add_speaker_notes(slide, notes or "Market sizing: TAM/SAM/SOM breakdown.")


def _build_business_model_slide(prs: Any, ctx: PitchContext, colors: dict[str, Any]) -> None:
    """Build business model / pricing slide."""
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches

    bm = ctx.spec.business_model
    if not bm or not bm.tiers:
        return

    slide = _create_dark_slide(prs, colors)
    _add_slide_heading(slide, "Business Model", colors)

    currency = ctx.spec.company.currency
    tier_count = len(bm.tiers)
    box_w = min(3.5, 11.0 / max(tier_count, 1))

    for i, tier in enumerate(bm.tiers):
        x = 0.8 + i * (box_w + 0.3)
        y = 2.0

        # Card background for each tier
        card_border = colors["accent"] if tier.highlighted else colors["muted"]
        card_h = min(2.3, CONTENT_BOTTOM - y + 0.1)
        _add_card(
            slide,
            Inches(x - 0.1),
            Inches(y - 0.1),
            Inches(box_w + 0.2),
            Inches(card_h),
            fill_color=colors.get("dark_text"),
            border_color=card_border,
        )

        name_color = colors["highlight"] if tier.highlighted else colors["white"]
        _add_text_box(
            slide,
            Inches(x),
            Inches(y),
            Inches(box_w),
            Inches(0.5),
            tier.name,
            font_size=22,
            bold=True,
            color=name_color,
            alignment=PP_ALIGN.CENTER,
        )

        price_str = _fmt_currency(tier.price, currency) if tier.price else "Free"
        _add_text_box(
            slide,
            Inches(x),
            Inches(y + 0.6),
            Inches(box_w),
            Inches(0.8),
            f"{price_str}/{tier.period}",
            font_size=28,
            bold=True,
            color=colors["accent"],
            alignment=PP_ALIGN.CENTER,
        )

        if tier.features:
            feature_h = _estimate_text_height(tier.features, box_w, 14)
            feature_box_h = min(card_h - 1.5, feature_h + 0.2)
            if feature_box_h > 0:
                _add_text_box(
                    slide,
                    Inches(x),
                    Inches(y + 1.5),
                    Inches(box_w),
                    Inches(feature_box_h),
                    tier.features,
                    font_size=14,
                    color=colors["muted"],
                    alignment=PP_ALIGN.CENTER,
                )

    notes = bm.speaker_notes
    _add_speaker_notes(slide, notes or f"Business model with {tier_count} pricing tiers.")


def _build_financials_slide(prs: Any, ctx: PitchContext, colors: dict[str, Any]) -> None:
    """Build financials slide with projections."""
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches

    fin = ctx.spec.financials
    if not fin or not fin.projections:
        return

    slide = _create_dark_slide(prs, colors)
    _add_slide_heading(slide, "Financial Projections", colors)

    currency = ctx.spec.company.currency
    col_count = len(fin.projections)
    gap = 0.3
    col_w = min(3.0, (11.0 - gap * (col_count - 1)) / max(col_count, 1))
    card_h = 2.4
    card_top = 2.2

    for i, proj in enumerate(fin.projections):
        x = 0.8 + i * (col_w + gap)
        # Card background for visual consistency
        _add_card(
            slide,
            Inches(x),
            Inches(card_top),
            Inches(col_w),
            Inches(card_h),
            fill_color=colors.get("dark_text"),
            border_color=colors["accent"],
        )
        # Year header
        _add_text_box(
            slide,
            Inches(x),
            Inches(card_top + 0.15),
            Inches(col_w),
            Inches(0.4),
            str(proj.year),
            font_size=16,
            bold=True,
            color=colors["muted"],
            alignment=PP_ALIGN.CENTER,
        )
        # Revenue value
        _add_text_box(
            slide,
            Inches(x),
            Inches(card_top + 0.55),
            Inches(col_w),
            Inches(0.6),
            _fmt_currency(proj.revenue, currency),
            font_size=24,
            bold=True,
            color=colors["success"],
            alignment=PP_ALIGN.CENTER,
        )
        # Revenue label
        _add_text_box(
            slide,
            Inches(x),
            Inches(card_top + 1.15),
            Inches(col_w),
            Inches(0.3),
            "Revenue",
            font_size=11,
            color=colors["muted"],
            alignment=PP_ALIGN.CENTER,
        )
        # Customer count
        _add_text_box(
            slide,
            Inches(x),
            Inches(card_top + 1.6),
            Inches(col_w),
            Inches(0.4),
            f"{proj.customers:,}",
            font_size=18,
            bold=True,
            color=colors["white"],
            alignment=PP_ALIGN.CENTER,
        )
        _add_text_box(
            slide,
            Inches(x),
            Inches(card_top + 2.0),
            Inches(col_w),
            Inches(0.3),
            "Customers",
            font_size=11,
            color=colors["muted"],
            alignment=PP_ALIGN.CENTER,
        )

    if fin.use_of_funds:
        y = 5.0
        _add_text_box(
            slide,
            Inches(0.8),
            Inches(y),
            Inches(11),
            Inches(0.4),
            "Use of Funds",
            font_size=16,
            bold=True,
            color=colors["accent"],
        )
        y += 0.45
        parts = [f"{f.category} ({f.percent}%)" for f in fin.use_of_funds]
        _add_text_box(
            slide,
            Inches(1.0),
            Inches(y),
            Inches(11),
            Inches(0.4),
            "  |  ".join(parts),
            font_size=13,
            color=colors["white"],
        )

    # Embed revenue chart if available
    revenue_chart = ctx.chart_paths.get("revenue")
    if revenue_chart and Path(revenue_chart).exists():
        try:
            slide.shapes.add_picture(
                str(revenue_chart), Inches(0.8), Inches(4.5), Inches(5), Inches(2.0)
            )
        except Exception as e:
            logger.debug(f"Failed to embed revenue chart: {e}")

    notes = fin.speaker_notes
    _add_speaker_notes(slide, notes or "Financial projections and use of funds breakdown.")


def _estimate_team_height(team: Any, scale: dict[str, int] | None = None) -> float:
    """Estimate total height needed for all team sections."""
    if scale is None:
        scale = {"name": 22, "role": 18, "bio": 14, "section": 20, "item": 16}
    height = 0.0
    for member in team.founders:
        height += 1.0 if member.bio else 0.7
    if team.advisors:
        height += 0.3 + 0.5  # gap + header
        height += len(team.advisors) * 0.5
    if team.key_hires:
        height += 0.3 + 0.5  # gap + header
        height += len(team.key_hires) * 0.5
    return height


def _render_team_member(
    slide: Any,
    region: ContentRegion,
    member: Any,
    colors: dict[str, Any],
    scale: dict[str, int],
    font_name: str | None,
) -> ContentRegion:
    """Render a single team member and return advanced region."""
    from pptx.util import Inches

    _add_text_box(
        slide,
        Inches(1.2),
        Inches(region.top),
        Inches(3),
        Inches(0.5),
        member.name,
        font_size=scale["name"],
        bold=True,
        color=colors["accent"],
        font_name=font_name,
    )
    _add_text_box(
        slide,
        Inches(4.5),
        Inches(region.top),
        Inches(2),
        Inches(0.5),
        member.role,
        font_size=scale["role"],
        color=colors["highlight"],
        font_name=font_name,
    )
    if member.bio:
        _add_text_box(
            slide,
            Inches(1.2),
            Inches(region.top + 0.4),
            Inches(10),
            Inches(0.4),
            member.bio,
            font_size=scale["bio"],
            color=colors["muted"],
            font_name=font_name,
        )
        return region.advance(1.0)
    return region.advance(0.7)


def _build_team_content(
    prs: Any,
    slide: Any,
    region: ContentRegion,
    colors: dict[str, Any],
    founders: list[Any],
    advisors: list[Any],
    key_hires: list[Any],
    scale: dict[str, int],
    title: str,
) -> None:
    """Render team content on a slide, creating continuation slides on overflow."""
    from pptx.util import Inches

    font_name = colors.get("font_family")

    overflow_founders: list[Any] = []
    overflow_advisors: list[Any] = []
    overflow_hires: list[Any] = []

    # Founders
    for idx, member in enumerate(founders):
        needed = 1.0 if member.bio else 0.7
        if not region.fits(needed):
            overflow_founders = founders[idx:]
            break
        region = _render_team_member(slide, region, member, colors, scale, font_name)

    # Advisors
    if not overflow_founders and advisors:
        header_h = 0.8  # gap + header
        if not region.fits(header_h + 0.5):
            overflow_advisors = advisors
        else:
            region = region.advance(0.3)
            _add_text_box(
                slide,
                Inches(0.8),
                Inches(region.top),
                Inches(11),
                Inches(0.5),
                "Advisors",
                font_size=scale["section"],
                bold=True,
                color=colors["accent"],
                font_name=font_name,
            )
            region = region.advance(0.5)
            for idx, advisor in enumerate(advisors):
                if not region.fits(0.5):
                    overflow_advisors = advisors[idx:]
                    break
                _add_text_box(
                    slide,
                    Inches(1.2),
                    Inches(region.top),
                    Inches(5),
                    Inches(0.4),
                    f"{advisor.name} \u2014 {advisor.role}",
                    font_size=scale["item"],
                    color=colors["dark_text"],
                    font_name=font_name,
                )
                region = region.advance(0.5)
    elif overflow_founders:
        overflow_advisors = advisors

    # Key Hires
    if not overflow_founders and not overflow_advisors and key_hires:
        header_h = 0.8
        if not region.fits(header_h + 0.5):
            overflow_hires = key_hires
        else:
            region = region.advance(0.3)
            _add_text_box(
                slide,
                Inches(0.8),
                Inches(region.top),
                Inches(11),
                Inches(0.5),
                "Key Hires",
                font_size=scale["section"],
                bold=True,
                color=colors["accent"],
                font_name=font_name,
            )
            region = region.advance(0.5)
            for idx, hire in enumerate(key_hires):
                if not region.fits(0.5):
                    overflow_hires = key_hires[idx:]
                    break
                timing = f" ({hire.timing})" if hire.timing else ""
                _add_text_box(
                    slide,
                    Inches(1.2),
                    Inches(region.top),
                    Inches(10),
                    Inches(0.4),
                    f"{hire.role}{timing}",
                    font_size=scale["item"],
                    color=colors["dark_text"],
                    font_name=font_name,
                )
                region = region.advance(0.5)
    elif overflow_founders or overflow_advisors:
        overflow_hires = key_hires

    # Create continuation slide if needed
    if overflow_founders or overflow_advisors or overflow_hires:
        logger.warning("Team slide: content overflow, creating continuation slide")
        cont_slide = _create_light_slide(prs, colors)
        cont_y = _add_slide_heading(
            cont_slide, f"{title} (continued)", colors, text_color=colors["dark_text"]
        )
        cont_region = ContentRegion(left=0.8, top=cont_y, width=11.0, bottom=CONTENT_BOTTOM)
        _build_team_content(
            prs,
            cont_slide,
            cont_region,
            colors,
            overflow_founders,
            overflow_advisors,
            overflow_hires,
            scale,
            title,
        )


def _build_team_slide(prs: Any, ctx: PitchContext, colors: dict[str, Any]) -> None:
    """Build team slide with light background, auto-scale, and auto-split."""
    team = ctx.spec.team
    if not team:
        return

    # Font scale options: normal -> compact -> minimum
    font_scales = [
        {"name": 22, "role": 18, "bio": 14, "section": 20, "item": 16},
        {"name": 18, "role": 16, "bio": 13, "section": 18, "item": 14},
        {"name": 16, "role": 14, "bio": 12, "section": 16, "item": 12},
    ]

    available = CONTENT_BOTTOM - 2.0  # heading takes ~2.0"
    scale = font_scales[0]
    for s in font_scales:
        total = _estimate_team_height(team, s)
        if total <= available:
            scale = s
            break
        scale = s  # use smallest if none fit

    slide = _create_light_slide(prs, colors)
    y = _add_slide_heading(slide, "Team", colors, text_color=colors["dark_text"])
    region = ContentRegion(left=0.8, top=y, width=11.0, bottom=CONTENT_BOTTOM)

    _build_team_content(
        prs,
        slide,
        region,
        colors,
        list(team.founders),
        list(team.advisors) if team.advisors else [],
        list(team.key_hires) if team.key_hires else [],
        scale,
        "Team",
    )

    notes = team.speaker_notes
    _add_speaker_notes(
        slide,
        notes or f"Team: {len(team.founders)} founders, {len(team.advisors)} advisors.",
    )


def _build_competition_slide(prs: Any, ctx: PitchContext, colors: dict[str, Any]) -> None:
    """Build competition slide with light background and table."""
    from pptx.util import Inches

    if not ctx.spec.competitors:
        return

    slide = _create_light_slide(prs, colors)
    _add_slide_heading(slide, "Competitive Landscape", colors, text_color=colors["dark_text"])

    headers = ["Competitor", "Strength", "Weakness"]
    rows = [
        [comp.name, comp.strength or "", comp.weakness or ""] for comp in ctx.spec.competitors[:6]
    ]
    _add_table(
        slide,
        Inches(0.8),
        Inches(2.2),
        Inches(11.5),
        headers,
        rows,
        colors,
        col_widths=[3.5, 4.0, 4.0],
        font_size=14,
    )  # returns (shape, LayoutResult) — unused here

    _add_speaker_notes(slide, f"{len(ctx.spec.competitors)} competitors analyzed.")


def _build_milestones_content(
    prs: Any,
    slide: Any,
    region: ContentRegion,
    colors: dict[str, Any],
    sections: list[tuple[str, list[str], Any, str, Any | None]],
    section_index: int,
    is_first_on_slide: bool,
) -> None:
    """Render milestone sections on a slide, creating continuation on overflow."""
    from pptx.util import Inches

    font_name = colors.get("font_family")

    for i in range(section_index, len(sections)):
        title, items, title_color, bullet_char, item_color = sections[i]
        needs_divider = not is_first_on_slide or i > section_index
        divider_h = 0.2 if needs_divider else 0.0
        header_h = 0.4
        item_spacing = 0.4

        if not region.fits(divider_h + header_h + item_spacing):
            # Can't fit even the header + one item — continue on next slide
            logger.warning("Milestones slide: content overflow, creating continuation slide")
            cont_slide = _create_dark_slide(prs, colors)
            cont_y = _add_slide_heading(cont_slide, "Milestones & Roadmap (continued)", colors)
            cont_region = ContentRegion(left=1.2, top=cont_y, width=10.0, bottom=CONTENT_BOTTOM)
            _build_milestones_content(prs, cont_slide, cont_region, colors, sections, i, True)
            return

        if needs_divider:
            _add_divider(slide, region.top + 0.05, colors)
            region = region.advance(divider_h)

        _add_text_box(
            slide,
            Inches(0.8),
            Inches(region.top),
            Inches(5),
            Inches(header_h),
            title,
            font_size=18,
            bold=True,
            color=title_color,
            font_name=font_name,
        )
        region = region.advance(header_h)

        lr = _add_bullet_list(
            slide,
            Inches(1.2),
            region.top,
            Inches(10),
            items,
            colors,
            font_size=15,
            spacing=item_spacing,
            bullet_char=bullet_char,
            color=item_color or colors["white"],
        )
        region = ContentRegion(
            left=region.left, top=lr.final_y, width=region.width, bottom=region.bottom
        )
        is_first_on_slide = False


def _build_milestones_slide(prs: Any, ctx: PitchContext, colors: dict[str, Any]) -> None:
    """Build milestones / roadmap slide with auto-split on overflow."""
    ms = ctx.spec.milestones
    if not ms:
        return

    slide = _create_dark_slide(prs, colors)
    y = _add_slide_heading(slide, "Milestones & Roadmap", colors)
    region = ContentRegion(left=1.2, top=y, width=10.0, bottom=CONTENT_BOTTOM)

    sections: list[tuple[str, list[str], Any, str, Any | None]] = []
    if ms.completed:
        sections.append(
            ("Completed \u2713", ms.completed, colors["success"], "\u2713", colors["muted"])
        )
    if ms.next_12_months:
        sections.append(("Next 12 Months", ms.next_12_months, colors["accent"], "\u2192", None))
    if ms.long_term:
        sections.append(
            ("Long Term Vision", ms.long_term, colors["highlight"], "\u25c6", colors["muted"])
        )

    _build_milestones_content(prs, slide, region, colors, sections, 0, True)

    notes = ms.speaker_notes
    _add_speaker_notes(slide, notes or "Milestones and roadmap.")


def _build_ask_slide(prs: Any, ctx: PitchContext, colors: dict[str, Any]) -> None:
    """Build the funding ask slide."""
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches

    if not ctx.spec.company.funding_ask:
        return

    slide = _create_dark_slide(prs, colors)
    _add_slide_heading(slide, "The Ask", colors)

    currency = ctx.spec.company.currency
    ask_str = _fmt_currency(ctx.spec.company.funding_ask, currency)
    stage = ctx.spec.company.stage.value.replace("_", " ").title()

    _add_text_box(
        slide,
        Inches(1),
        Inches(2.5),
        Inches(11),
        Inches(1.5),
        ask_str,
        font_size=64,
        bold=True,
        color=colors["highlight"],
        alignment=PP_ALIGN.CENTER,
    )
    _add_text_box(
        slide,
        Inches(1),
        Inches(4.0),
        Inches(11),
        Inches(0.5),
        f"{stage} Round",
        font_size=24,
        color=colors["accent"],
        alignment=PP_ALIGN.CENTER,
    )

    if ctx.spec.company.runway_months:
        _add_text_box(
            slide,
            Inches(1),
            Inches(5.0),
            Inches(11),
            Inches(0.5),
            f"{ctx.spec.company.runway_months} months runway",
            font_size=18,
            color=colors["muted"],
            alignment=PP_ALIGN.CENTER,
        )

    if ctx.spec.financials and ctx.spec.financials.use_of_funds:
        y = 5.8
        parts = [f"{f.category} ({f.percent}%)" for f in ctx.spec.financials.use_of_funds]
        _add_text_box(
            slide,
            Inches(1),
            Inches(y),
            Inches(11),
            Inches(0.5),
            "  |  ".join(parts),
            font_size=14,
            color=colors["white"],
            alignment=PP_ALIGN.CENTER,
        )

    # Embed funds chart if available
    funds_chart = ctx.chart_paths.get("funds")
    if funds_chart and Path(funds_chart).exists():
        try:
            slide.shapes.add_picture(
                str(funds_chart), Inches(8), Inches(4.5), Inches(4), Inches(2.5)
            )
        except Exception as e:
            logger.debug(f"Failed to embed funds chart: {e}")

    _add_speaker_notes(slide, f"Raising {ask_str} at {stage} stage.")


def _build_closing_slide(prs: Any, ctx: PitchContext, colors: dict[str, Any]) -> None:
    """Build the closing / thank you slide."""
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches

    slide = _create_dark_slide(prs, colors)
    _add_text_box(
        slide,
        Inches(1),
        Inches(2.5),
        Inches(11),
        Inches(1.5),
        "Thank You",
        font_size=48,
        bold=True,
        color=colors["white"],
        alignment=PP_ALIGN.CENTER,
    )
    _add_text_box(
        slide,
        Inches(1),
        Inches(4.2),
        Inches(11),
        Inches(0.8),
        ctx.spec.company.name,
        font_size=28,
        color=colors["accent"],
        alignment=PP_ALIGN.CENTER,
    )
    if ctx.spec.company.tagline:
        _add_text_box(
            slide,
            Inches(1),
            Inches(5.0),
            Inches(11),
            Inches(0.5),
            ctx.spec.company.tagline,
            font_size=18,
            color=colors["muted"],
            alignment=PP_ALIGN.CENTER,
        )

    notes = ctx.spec.company.speaker_notes if ctx.spec.company else None
    _add_speaker_notes(slide, notes or "Closing slide. Open for questions.")


def _build_extra_slide(
    prs: Any, ctx: PitchContext, colors: dict[str, Any], extra: ExtraSlide
) -> None:
    """Build a user-defined extra slide based on its layout."""
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches

    use_light = extra.theme == "light"
    slide = _create_light_slide(prs, colors) if use_light else _create_dark_slide(prs, colors)
    heading_color = colors["dark_text"] if use_light else colors["white"]
    text_color = colors["dark_text"] if use_light else colors["white"]
    y = _add_slide_heading(slide, extra.title, colors, text_color=heading_color)

    font_name = colors.get("font_family")

    if extra.layout == ExtraSlideLayout.BULLETS:
        _add_bullet_list(
            slide,
            Inches(1.2),
            y,
            Inches(10),
            extra.items,
            colors,
            font_size=20,
            spacing=0.7,
            color=text_color,
        )

    elif extra.layout == ExtraSlideLayout.STATS:
        parsed: list[tuple[str, str]] = []
        for item in extra.items:
            if "|" in item:
                val, lbl = item.split("|", 1)
                parsed.append((val.strip(), lbl.strip()))
            else:
                parsed.append((item, ""))
        # Cap items when boxes would be too narrow to read
        max_items = min(6, len(parsed))
        if max_items > 0:
            box_w = 11.0 / max_items
            if box_w < 2.0 and max_items > 4:
                logger.warning(
                    f"Stats layout: capping from {max_items} to 4 items (boxes too narrow)"
                )
                max_items = 4
        parsed = parsed[:max_items]
        _add_columns(
            slide,
            Inches(2.0),
            parsed,
            value_color=colors["accent"],
            label_color=colors["muted"],
            alignment=PP_ALIGN.CENTER,
        )

    elif extra.layout == ExtraSlideLayout.CARDS:
        # Multi-column card grid
        items = extra.items
        col_count = min(3, len(items)) if items else 1
        margin = 1.0
        gap = 0.3
        card_w = (11.0 - gap * (col_count - 1)) / col_count
        card_h = 0.8
        row_gap = 0.3
        for idx, item in enumerate(items):
            col = idx % col_count
            row = idx // col_count
            x = margin + col * (card_w + gap)
            cy = y + row * (card_h + row_gap)
            if cy + card_h > CONTENT_BOTTOM:
                remaining = len(items) - idx
                logger.warning(f"Cards layout: truncated {remaining} cards due to overflow")
                break
            _add_card(
                slide,
                Inches(x),
                Inches(cy),
                Inches(card_w),
                Inches(card_h),
                fill_color=colors.get("dark_text"),
                border_color=colors["accent"],
            )
            _add_rich_text_box(
                slide,
                Inches(x + 0.3),
                Inches(cy + 0.1),
                Inches(card_w - 0.6),
                Inches(0.6),
                item,
                font_size=18,
                color=colors["white"] if not use_light else colors["white"],
                font_name=font_name,
            )

    elif extra.layout == ExtraSlideLayout.TABLE:
        # Parse items as pipe-separated: first item = headers, rest = rows
        if extra.items:
            headers = [h.strip() for h in extra.items[0].split("|")]
            rows = [[c.strip() for c in item.split("|")] for item in extra.items[1:]]
            _add_table(
                slide,
                Inches(0.8),
                Inches(y),
                Inches(11.5),
                headers,
                rows,
                colors,
                font_size=14,
            )  # returns (shape, LayoutResult) — unused here

    elif extra.layout == ExtraSlideLayout.CALLOUT:
        # First item as callout box, rest as supporting bullets
        if extra.items:
            _add_callout_box(
                slide,
                Inches(0.8),
                Inches(y),
                Inches(11.5),
                extra.items[0],
                colors,
                font_size=24,
            )
            if len(extra.items) > 1:
                _add_bullet_list(
                    slide,
                    Inches(1.2),
                    y + 1.5,
                    Inches(10),
                    extra.items[1:],
                    colors,
                    font_size=18,
                    spacing=0.6,
                    color=text_color,
                )

    elif extra.layout == ExtraSlideLayout.CUSTOM:
        # Handled by plugin system in pptx_gen.py
        pass

    elif extra.layout == ExtraSlideLayout.IMAGE:
        if extra.image_path:
            try:
                img_h = min(5.0, CONTENT_BOTTOM - y)
                slide.shapes.add_picture(
                    extra.image_path,
                    Inches(1.5),
                    Inches(y),
                    Inches(10),
                    Inches(img_h),
                )
            except Exception as e:
                logger.warning(f"Failed to embed image {extra.image_path}: {e}")

    notes = extra.speaker_notes
    _add_speaker_notes(slide, notes or f"Extra slide: {extra.title}")
