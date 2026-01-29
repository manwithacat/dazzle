"""
Plugin loader for custom pitch slide builders.

Discovers project-local slide plugins from a ``pitch_slides/`` directory
and registers their builder functions for use with ``ExtraSlideLayout.CUSTOM``.
"""

from __future__ import annotations

import importlib.util
import inspect
import logging
from collections.abc import Callable
from pathlib import Path

logger = logging.getLogger(__name__)


class PluginRegistry:
    """Registry of custom slide builder functions."""

    def __init__(self) -> None:
        self._builders: dict[str, Callable[..., None]] = {}

    def register(self, name: str, builder: Callable[..., None]) -> None:
        """Register a builder function by name."""
        self._builders[name] = builder
        logger.debug(f"Registered slide builder: {name}")

    def get(self, name: str) -> Callable[..., None] | None:
        """Get a builder function by name."""
        return self._builders.get(name)

    def list_builders(self) -> list[str]:
        """List all registered builder names."""
        return sorted(self._builders.keys())


def discover_plugins(project_root: Path) -> PluginRegistry:
    """Discover slide builder plugins from ``pitch_slides/`` directory.

    Scans for Python files containing ``build_*_slide`` functions with 4+
    parameters (prs, ctx, colors, extra).

    Args:
        project_root: Project root directory containing ``pitch_slides/``.

    Returns:
        PluginRegistry with discovered builders.
    """
    registry = PluginRegistry()
    plugin_dir = project_root / "pitch_slides"

    if not plugin_dir.is_dir():
        return registry

    for py_file in sorted(plugin_dir.glob("*.py")):
        if py_file.name.startswith("_"):
            continue
        try:
            spec = importlib.util.spec_from_file_location(f"pitch_slides.{py_file.stem}", py_file)
            if spec is None or spec.loader is None:
                continue
            module = importlib.util.module_from_spec(spec)
            spec.loader.exec_module(module)

            for attr_name in dir(module):
                if not attr_name.startswith("build_") or not attr_name.endswith("_slide"):
                    continue
                fn = getattr(module, attr_name)
                if not callable(fn):
                    continue
                sig = inspect.signature(fn)
                if len(sig.parameters) < 4:
                    continue
                # Register as e.g. "my_custom" from "build_my_custom_slide"
                builder_name = attr_name[len("build_") : -len("_slide")]
                registry.register(builder_name, fn)

        except Exception as e:
            logger.warning(f"Failed to load plugin {py_file.name}: {e}")

    return registry


def generate_plugin_readme(project_root: Path) -> Path:
    """Generate a README.md documenting the slide builder plugin API.

    Args:
        project_root: Project root where ``pitch_slides/`` lives.

    Returns:
        Path to the generated README.md.
    """
    plugin_dir = project_root / "pitch_slides"
    plugin_dir.mkdir(parents=True, exist_ok=True)

    readme_path = plugin_dir / "README.md"
    content = _readme_content()
    readme_path.write_text(content)
    return readme_path


def _readme_content() -> str:
    return """\
# Pitch Slide Plugins

<!-- AUTO-GENERATED — do not edit manually -->

Place Python files in this directory to add custom slide builders.

## Builder API

Each builder function must:

1. Be named `build_<name>_slide` (e.g. `build_roadmap_slide`)
2. Accept exactly 4 positional parameters:

```python
def build_roadmap_slide(prs, ctx, colors, extra):
    \"\"\"Build a custom roadmap slide.

    Args:
        prs: python-pptx Presentation object
        ctx: PitchContext with spec and DSL data
        colors: dict of resolved RGBColor values
        extra: ExtraSlide with title, items, data, etc.
    \"\"\"
    from pptx.util import Inches
    from dazzle.pitch.generators.pptx_primitives import (
        _create_dark_slide,
        _add_slide_heading,
        _add_bullet_list,
    )

    slide = _create_dark_slide(prs, colors)
    y = _add_slide_heading(slide, extra.title, colors)
    _add_bullet_list(slide, Inches(1.2), y, Inches(10), extra.items, colors)
```

## Using in pitchspec.yaml

```yaml
extra_slides:
  - title: "Roadmap"
    layout: custom
    builder: roadmap
    data:
      timeline: "2025-2027"
    items:
      - "Phase 1: MVP"
      - "Phase 2: Scale"
```

## Available Primitives

Import from `dazzle.pitch.generators.pptx_primitives`:

- `_create_dark_slide`, `_create_light_slide`
- `_add_slide_heading`, `_add_text_box`, `_add_rich_text_box`
- `_add_bullet_list`, `_add_table`, `_add_card`
- `_add_stat_box`, `_add_columns`, `_add_callout_box`
- `_add_divider`, `_add_speaker_notes`
- `LayoutResult`, `SLIDE_WIDTH`, `SLIDE_HEIGHT`, `CONTENT_TOP`, `CONTENT_BOTTOM`
"""
