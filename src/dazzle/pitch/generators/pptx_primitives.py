"""
PPTX visual primitives and helper functions.

Low-level shape helpers used by slide builders.
"""

from __future__ import annotations

import logging
import math
import re
from dataclasses import dataclass, replace
from typing import TYPE_CHECKING, Any

if TYPE_CHECKING:
    from pptx.dml.color import RGBColor
    from pptx.presentation import Presentation

    from dazzle.pitch.ir import BrandColors

logger = logging.getLogger(__name__)

# Layout constants (inches)
SLIDE_WIDTH = 13.333
SLIDE_HEIGHT = 7.5
CONTENT_TOP = 2.0  # after heading + accent bar
CONTENT_BOTTOM = 7.0  # 0.5" bottom margin
CONTENT_LEFT = 0.8
CONTENT_RIGHT = 12.5


@dataclass
class ContentRegion:
    """Vertical layout cursor with bounds tracking."""

    left: float  # inches
    top: float  # current y position (inches)
    width: float  # available width (inches)
    bottom: float  # hard stop y (inches), typically CONTENT_BOTTOM

    @property
    def remaining(self) -> float:
        return max(0.0, self.bottom - self.top)

    def fits(self, height: float) -> bool:
        return self.top + height <= self.bottom

    def advance(self, height: float) -> ContentRegion:
        return replace(self, top=self.top + height)


def _estimate_text_height(
    text: str, width_inches: float, font_size: int, *, bold: bool = False
) -> float:
    """Estimate rendered text height in inches (conservative/rounds up)."""
    char_width = font_size * 0.6  # points per character (approximate)
    if bold:
        char_width *= 1.1
    chars_per_line = max(1, int(width_inches * 72 / char_width))
    lines = 0
    for paragraph in text.split("\n"):
        lines += max(1, math.ceil(len(paragraph) / chars_per_line))
    line_height = font_size * 1.4 / 72  # inches
    return lines * line_height


@dataclass
class LayoutResult:
    """Result of a layout operation with overflow tracking."""

    final_y: float
    items_rendered: int
    items_truncated: int
    overflow: bool


def _resolve_colors(brand: BrandColors) -> dict[str, Any]:
    """Convert brand hex colors to RGBColor objects."""
    from pptx.dml.color import RGBColor

    def hex_to_rgb(hex_color: str) -> RGBColor:
        h = hex_color.lstrip("#")
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))  # type: ignore[no-untyped-call]

    return {
        "primary": hex_to_rgb(brand.primary),
        "accent": hex_to_rgb(brand.accent),
        "highlight": hex_to_rgb(brand.highlight),
        "success": hex_to_rgb(brand.success),
        "light": hex_to_rgb(brand.light),
        "white": RGBColor(0xFF, 0xFF, 0xFF),  # type: ignore[no-untyped-call]
        "dark_text": hex_to_rgb(brand.primary),
        "muted": RGBColor(0x99, 0x99, 0x99),  # type: ignore[no-untyped-call]
        "font_family": brand.font_family,
    }


def _add_text_box(
    slide: Any,
    left: int,
    top: int,
    width: int,
    height: int,
    text: str,
    *,
    font_size: int = 18,
    bold: bool = False,
    color: RGBColor | None = None,
    alignment: int | None = None,
    font_name: str | None = None,
) -> Any:
    """Add a text box to a slide."""
    from pptx.util import Pt

    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    if color:
        p.font.color.rgb = color
    if alignment is not None:
        p.alignment = alignment
    if font_name:
        p.font.name = font_name
    return txbox


def _add_rich_text_box(
    slide: Any,
    left: int,
    top: int,
    width: int,
    height: int,
    text: str,
    *,
    font_size: int = 18,
    color: RGBColor | None = None,
    alignment: int | None = None,
    font_name: str | None = None,
) -> Any:
    """Add a text box that parses **bold** markers into bold runs."""
    from pptx.util import Pt

    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if alignment is not None:
        p.alignment = alignment

    parts = re.split(r"(\*\*.*?\*\*)", text)
    for part in parts:
        if part.startswith("**") and part.endswith("**"):
            run = p.add_run()
            run.text = part[2:-2]
            run.font.bold = True
            run.font.size = Pt(font_size)
            if color:
                run.font.color.rgb = color
            if font_name:
                run.font.name = font_name
        elif part:
            run = p.add_run()
            run.text = part
            run.font.size = Pt(font_size)
            run.font.bold = False
            if color:
                run.font.color.rgb = color
            if font_name:
                run.font.name = font_name
    return txbox


def _create_dark_slide(prs: Presentation, colors: dict[str, Any]) -> Any:
    """Create a slide with dark background."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = colors["primary"]
    return slide


def _create_light_slide(prs: Presentation, colors: dict[str, Any]) -> Any:
    """Create a slide with light background."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = colors["light"]
    return slide


def _add_speaker_notes(slide: Any, text: str) -> None:
    """Add speaker notes to a slide."""
    notes_slide = slide.notes_slide
    if notes_slide.notes_text_frame is not None:
        notes_slide.notes_text_frame.text = text


def _fmt_currency(amount: int, currency: str = "GBP") -> str:
    """Format an amount as currency."""
    symbols = {"GBP": "\u00a3", "USD": "$", "EUR": "\u20ac"}
    symbol = symbols.get(currency, currency + " ")
    if amount >= 1_000_000_000:
        return f"{symbol}{amount / 1_000_000_000:.1f}B"
    elif amount >= 1_000_000:
        return f"{symbol}{amount / 1_000_000:.1f}M"
    elif amount >= 1_000:
        return f"{symbol}{amount / 1_000:.0f}K"
    return f"{symbol}{amount:,}"


def _add_card(
    slide: Any,
    left: int,
    top: int,
    width: int,
    height: int,
    *,
    fill_color: RGBColor | None = None,
    border_color: RGBColor | None = None,
) -> Any:
    """Add a rounded rectangle card shape to a slide."""
    from pptx.util import Pt

    shape = slide.shapes.add_shape(
        5,  # MSO_SHAPE.ROUNDED_RECTANGLE
        left,
        top,
        width,
        height,
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def _add_stat_box(
    slide: Any,
    left: int,
    top: int,
    width: int,
    value: str,
    label: str,
    *,
    value_color: RGBColor | None = None,
    label_color: RGBColor | None = None,
    alignment: int | None = None,
) -> None:
    """Add a big-number + caption stat box."""
    from pptx.util import Inches

    # Convert width to inches for estimation
    width_in = width / Inches(1) if not isinstance(width, float) else width

    # Auto-scale value font to fit within the box width
    value_font = 48
    while _estimate_text_height(value, width_in, value_font, bold=True) > 1.0 and value_font > 24:
        value_font -= 4

    # Auto-scale label font to fit
    label_font = 16
    while _estimate_text_height(label, width_in, label_font) > 0.5 and label_font > 11:
        label_font -= 1

    _add_text_box(
        slide,
        left,
        top,
        width,
        Inches(1),
        value,
        font_size=value_font,
        bold=True,
        color=value_color,
        alignment=alignment,
    )
    _add_text_box(
        slide,
        left,
        top + Inches(1),
        width,
        Inches(0.5),
        label,
        font_size=label_font,
        color=label_color,
        alignment=alignment,
    )


def _add_columns(
    slide: Any,
    top: int,
    items: list[tuple[str, str]],
    *,
    value_color: RGBColor | None = None,
    label_color: RGBColor | None = None,
    alignment: int | None = None,
) -> None:
    """Auto-layout N stat boxes across slide width."""
    from pptx.util import Inches

    n = len(items)
    if n == 0:
        return
    box_w = min(3.0, 11.0 / max(n, 1))
    x_start = 1.0
    for i, (value, label) in enumerate(items[:6]):
        x = x_start + i * (box_w + 0.2)
        _add_stat_box(
            slide,
            Inches(x),
            top,
            Inches(box_w),
            value,
            label,
            value_color=value_color,
            label_color=label_color,
            alignment=alignment,
        )


def _add_slide_heading(
    slide: Any, text: str, colors: dict[str, Any], *, text_color: RGBColor | None = None
) -> float:
    """Add a title heading with accent bar underneath. Returns y position after heading."""
    from pptx.util import Inches, Pt

    heading_color = text_color or colors["white"]
    font_name = colors.get("font_family")

    # Auto-scale font to fit single line within the 11" heading width.
    # At 36pt bold, ~46 chars fit on one line. Scale down for longer titles.
    font_size = 36
    max_box_height = 0.85  # inches available above accent bar
    estimated_h = _estimate_text_height(text, 11.0, font_size, bold=True)
    while estimated_h > max_box_height and font_size > 20:
        font_size -= 2
        estimated_h = _estimate_text_height(text, 11.0, font_size, bold=True)

    _add_text_box(
        slide,
        Inches(0.8),
        Inches(0.5),
        Inches(11),
        Inches(1),
        text,
        font_size=font_size,
        bold=True,
        color=heading_color,
        font_name=font_name,
    )
    # Accent bar
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(0.8),
        Inches(1.35),
        Inches(2.5),
        Inches(0.06),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = colors["accent"]
    bar.line.fill.background()
    _ = Pt  # suppress unused import warning
    return 2.0


def _add_bullet_list(
    slide: Any,
    left: int | float,
    top: float,
    width: int | float,
    items: list[str],
    colors: dict[str, Any],
    *,
    font_size: int = 18,
    spacing: float = 0.6,
    bullet_char: str = "\u2022",
    color: RGBColor | None = None,
) -> LayoutResult:
    """Render N bullet items with bounds checking. Returns LayoutResult."""
    from pptx.util import Inches

    text_color = color or colors["white"]
    font_name = colors.get("font_family")
    y = top if isinstance(top, float) else top / 914400  # convert EMU to inches if needed
    # Normalize to float inches
    if not isinstance(y, float):
        y = float(y)
    left_emu = left if not isinstance(left, float) else Inches(left)
    width_emu = width if not isinstance(width, float) else Inches(width)
    rendered = 0
    for item in items:
        if y + spacing > CONTENT_BOTTOM:
            truncated = len(items) - rendered
            logger.warning(f"Bullet list overflow: truncated {truncated} items")
            return LayoutResult(
                final_y=y, items_rendered=rendered, items_truncated=truncated, overflow=True
            )
        _add_text_box(
            slide,
            left_emu,
            Inches(y),
            width_emu,
            Inches(spacing),
            f"{bullet_char} {item}",
            font_size=font_size,
            color=text_color,
            font_name=font_name,
        )
        y += spacing
        rendered += 1
    return LayoutResult(final_y=y, items_rendered=rendered, items_truncated=0, overflow=False)


def _add_table(
    slide: Any,
    left: int,
    top: int,
    width: int,
    headers: list[str],
    rows: list[list[str]],
    colors: dict[str, Any],
    *,
    col_widths: list[float] | None = None,
    font_size: int = 14,
) -> tuple[Any, LayoutResult]:
    """Create a table shape with styled header row and alternating row colors.

    Returns (table_shape, LayoutResult) tuple.
    """
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    col_count = len(headers)
    if col_count == 0:
        return None, LayoutResult(final_y=0, items_rendered=0, items_truncated=0, overflow=False)

    # Bounds check: truncate rows that won't fit
    top_inches = top if isinstance(top, float) else float(top / 914400)
    row_height = 0.4
    max_rows_fit = max(1, int((CONTENT_BOTTOM - top_inches) / row_height) - 1)
    actual_rows = rows[:max_rows_fit]
    truncated = len(rows) - len(actual_rows)
    if truncated > 0:
        logger.warning(f"Table overflow: truncated {truncated} rows")

    row_count = len(actual_rows) + 1  # +1 for header

    table_shape = slide.shapes.add_table(
        row_count, col_count, left, top, width, Inches(0.4 * row_count)
    )
    table = table_shape.table

    # Set column widths
    if col_widths:
        for i, w in enumerate(col_widths[:col_count]):
            table.columns[i].width = Inches(w)

    font_name = colors.get("font_family")

    # Header row
    for ci, header in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = header
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(font_size)
            paragraph.font.bold = True
            paragraph.font.color.rgb = colors["white"]
            if font_name:
                paragraph.font.name = font_name
        cell.fill.solid()
        cell.fill.fore_color.rgb = colors["accent"]

    # Data rows
    for ri, row in enumerate(actual_rows):
        for ci, value in enumerate(row[:col_count]):
            cell = table.cell(ri + 1, ci)
            cell.text = value
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.color.rgb = colors["dark_text"]
                if font_name:
                    paragraph.font.name = font_name
            # Alternating row colors
            cell.fill.solid()
            if ri % 2 == 0:
                cell.fill.fore_color.rgb = colors["light"]
            else:
                cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)  # type: ignore[no-untyped-call]

    final_y = top_inches + row_height * row_count
    lr = LayoutResult(
        final_y=final_y,
        items_rendered=len(actual_rows),
        items_truncated=truncated,
        overflow=truncated > 0,
    )
    return table_shape, lr


def _add_callout_box(
    slide: Any,
    left: int,
    top: int,
    width: int,
    text: str,
    colors: dict[str, Any],
    *,
    font_size: int = 24,
) -> Any:
    """Dark rounded rect with accent left border and bold text inside."""
    from pptx.util import Inches, Pt

    # Main box
    box_height = Inches(1.2)
    shape = slide.shapes.add_shape(
        5,  # MSO_SHAPE.ROUNDED_RECTANGLE
        left,
        top,
        width,
        box_height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = colors["primary"]
    shape.line.fill.background()

    # Accent left border
    border = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        left,
        top,
        Inches(0.08),
        box_height,
    )
    border.fill.solid()
    border.fill.fore_color.rgb = colors["accent"]
    border.line.fill.background()

    # Text
    font_name = colors.get("font_family")
    _add_text_box(
        slide,
        left + Inches(0.3),
        top + Inches(0.2),
        width - Inches(0.6),
        box_height - Inches(0.4),
        text,
        font_size=font_size,
        bold=True,
        color=colors["white"],
        font_name=font_name,
    )
    _ = Pt  # suppress unused
    return shape


def _add_timeline(
    slide: Any,
    top: float,
    items: list[str],
    colors: dict[str, Any],
    *,
    width: float = 11.0,
) -> None:
    """Render a horizontal timeline with dots, connecting line, and labels.

    Each item is a ``"date|description"`` pipe-separated string.
    """
    from pptx.util import Inches, Pt

    n = len(items)
    if n == 0:
        return

    font_name = colors.get("font_family")
    left = 1.0
    line_y = top + 0.4

    # Connecting line
    line = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(left),
        Inches(line_y),
        Inches(width),
        Inches(0.04),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = colors["accent"]
    line.line.fill.background()

    spacing = width / max(n, 1)
    dot_size = 0.22

    for i, item in enumerate(items):
        parts = item.split("|", 1)
        date_text = parts[0].strip()
        desc_text = parts[1].strip() if len(parts) > 1 else ""

        cx = left + i * spacing + spacing / 2
        dot_left = cx - dot_size / 2

        # Dot
        dot = slide.shapes.add_shape(
            9,  # MSO_SHAPE.OVAL
            Inches(dot_left),
            Inches(line_y - dot_size / 2 + 0.02),
            Inches(dot_size),
            Inches(dot_size),
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = colors["highlight"]
        dot.line.fill.background()

        # Date label (above line)
        _add_text_box(
            slide,
            Inches(cx - spacing / 2),
            Inches(top - 0.1),
            Inches(spacing),
            Inches(0.4),
            date_text,
            font_size=12,
            bold=True,
            color=colors["accent"],
            alignment=1,  # PP_ALIGN.CENTER
            font_name=font_name,
        )

        # Description (below line)
        _add_text_box(
            slide,
            Inches(cx - spacing / 2),
            Inches(line_y + 0.25),
            Inches(spacing),
            Inches(0.8),
            desc_text,
            font_size=11,
            color=colors["white"],
            alignment=1,  # PP_ALIGN.CENTER
            font_name=font_name,
        )

    _ = Pt  # suppress unused


def _add_divider(slide: Any, y: float, colors: dict[str, Any]) -> Any:
    """Thin horizontal accent line."""
    from pptx.util import Inches

    bar = slide.shapes.add_shape(1, Inches(0.8), Inches(y), Inches(11.5), Inches(0.03))
    bar.fill.solid()
    bar.fill.fore_color.rgb = colors["accent"]
    bar.line.fill.background()
    return bar
