"""Tests for pitch deck generators."""

from __future__ import annotations

from pathlib import Path

import pytest

from dazzle.pitch.extractor import PitchContext
from dazzle.pitch.ir import (
    BusinessModelSpec,
    CompanySpec,
    Competitor,
    ExtraSlide,
    ExtraSlideLayout,
    FinancialsSpec,
    FundAllocation,
    FundingStage,
    MarketSize,
    MarketSpec,
    MilestonesSpec,
    PitchSpec,
    PricingTier,
    ProblemSpec,
    SolutionSpec,
    TeamMember,
    TeamSpec,
    YearProjection,
)


def _make_full_context() -> PitchContext:
    """Create a fully-populated PitchContext for testing."""
    spec = PitchSpec(
        company=CompanySpec(
            name="TestCo",
            tagline="Testing pitch generation",
            stage=FundingStage.SEED,
            funding_ask=500000,
            runway_months=18,
        ),
        problem=ProblemSpec(
            headline="Testing is hard",
            points=["Point 1", "Point 2"],
            market_failure=["No good tools"],
        ),
        solution=SolutionSpec(
            headline="We make it easy",
            how_it_works=["Step 1", "Step 2"],
            value_props=["Fast", "Reliable"],
        ),
        market=MarketSpec(
            tam=MarketSize(value=10_000_000_000, label="TAM"),
            sam=MarketSize(value=1_000_000_000, label="SAM"),
            som=MarketSize(value=50_000_000, label="SOM"),
        ),
        business_model=BusinessModelSpec(
            tiers=[
                PricingTier(name="Free", price=0),
                PricingTier(name="Pro", price=49, highlighted=True),
            ]
        ),
        financials=FinancialsSpec(
            projections=[
                YearProjection(year=2025, customers=100, revenue=50000, costs=200000),
                YearProjection(year=2026, customers=500, revenue=300000, costs=350000),
            ],
            use_of_funds=[
                FundAllocation(category="Eng", percent=60),
                FundAllocation(category="Sales", percent=40),
            ],
        ),
        team=TeamSpec(
            founders=[TeamMember(name="Alice", role="CEO", bio="Expert")],
        ),
        competitors=[Competitor(name="BigCo", strength="Large", weakness="Slow")],
        milestones=MilestonesSpec(
            completed=["MVP"],
            next_12_months=["Launch"],
        ),
    )
    ctx = PitchContext(
        spec=spec,
        entities=["Task", "Project"],
        surfaces=["task_list", "project_detail"],
        personas=[{"id": "admin", "label": "Admin", "description": "System admin"}],
    )
    return ctx


class TestPptxGenerator:
    def test_generate_pptx(self, tmp_path: Path):
        """Test PPTX generation with full context."""
        try:
            import pptx  # noqa: F401
        except ImportError:
            pytest.skip("python-pptx not installed")

        from dazzle.pitch.generators.pptx_gen import generate_pptx

        ctx = _make_full_context()
        output = tmp_path / "test_deck.pptx"
        result = generate_pptx(ctx, output)

        assert result.success
        assert result.output_path == output
        assert output.exists()
        assert result.slide_count > 0

    def test_generate_pptx_minimal(self, tmp_path: Path):
        """Test PPTX generation with minimal context."""
        try:
            import pptx  # noqa: F401
        except ImportError:
            pytest.skip("python-pptx not installed")

        from dazzle.pitch.generators.pptx_gen import generate_pptx

        spec = PitchSpec()
        ctx = PitchContext(spec=spec)
        output = tmp_path / "minimal.pptx"
        result = generate_pptx(ctx, output)

        assert result.success
        # Should have at least title + closing
        assert result.slide_count >= 2

    def test_fmt_currency(self):
        from dazzle.pitch.generators.pptx_gen import _fmt_currency

        assert _fmt_currency(500000, "GBP") == "£500K"
        assert _fmt_currency(1500000, "USD") == "$1.5M"
        assert _fmt_currency(10_000_000_000, "EUR") == "€10.0B"
        assert _fmt_currency(42, "GBP") == "£42"


class TestVisualPrimitives:
    def test_add_stat_box(self, tmp_path: Path):
        try:
            import pptx  # noqa: F401
        except ImportError:
            pytest.skip("python-pptx not installed")

        from pptx import Presentation
        from pptx.util import Inches

        from dazzle.pitch.generators.pptx_gen import _add_stat_box, _resolve_colors
        from dazzle.pitch.ir import BrandColors

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        colors = _resolve_colors(BrandColors())
        _add_stat_box(
            slide,
            Inches(1),
            Inches(1),
            Inches(3),
            "42",
            "Users",
            value_color=colors["accent"],
            label_color=colors["muted"],
        )
        # Should have added 2 text boxes
        assert len(slide.shapes) >= 2

    def test_add_columns(self, tmp_path: Path):
        try:
            import pptx  # noqa: F401
        except ImportError:
            pytest.skip("python-pptx not installed")

        from pptx import Presentation
        from pptx.util import Inches

        from dazzle.pitch.generators.pptx_gen import _add_columns, _resolve_colors
        from dazzle.pitch.ir import BrandColors

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        colors = _resolve_colors(BrandColors())
        _add_columns(
            slide,
            Inches(2),
            [("10", "A"), ("20", "B"), ("30", "C")],
            value_color=colors["accent"],
            label_color=colors["muted"],
        )
        # 3 columns * 2 shapes each = 6
        assert len(slide.shapes) >= 6

    def test_add_rich_text_box_bold(self, tmp_path: Path):
        try:
            import pptx  # noqa: F401
        except ImportError:
            pytest.skip("python-pptx not installed")

        from pptx import Presentation
        from pptx.util import Inches

        from dazzle.pitch.generators.pptx_gen import _add_rich_text_box

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        txbox = _add_rich_text_box(
            slide,
            Inches(1),
            Inches(1),
            Inches(5),
            Inches(1),
            "Hello **world** today",
        )
        tf = txbox.text_frame
        runs = tf.paragraphs[0].runs
        assert len(runs) == 3
        assert runs[0].text == "Hello "
        assert not runs[0].font.bold
        assert runs[1].text == "world"
        assert runs[1].font.bold
        assert runs[2].text == " today"

    def test_add_card(self, tmp_path: Path):
        try:
            import pptx  # noqa: F401
        except ImportError:
            pytest.skip("python-pptx not installed")

        from pptx import Presentation
        from pptx.util import Inches

        from dazzle.pitch.generators.pptx_gen import _add_card, _resolve_colors
        from dazzle.pitch.ir import BrandColors

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        colors = _resolve_colors(BrandColors())
        shape = _add_card(
            slide,
            Inches(1),
            Inches(1),
            Inches(4),
            Inches(2),
            fill_color=colors["primary"],
            border_color=colors["accent"],
        )
        assert shape is not None


class TestExtraSlides:
    def test_extra_slide_bullets(self, tmp_path: Path):
        try:
            import pptx  # noqa: F401
        except ImportError:
            pytest.skip("python-pptx not installed")

        from dazzle.pitch.generators.pptx_gen import generate_pptx

        spec = PitchSpec(
            extra_slides=[
                ExtraSlide(title="Case Study", items=["Point 1", "Point 2"]),
            ]
        )
        ctx = PitchContext(spec=spec)
        output = tmp_path / "extra.pptx"
        result = generate_pptx(ctx, output)
        assert result.success
        # title + extra + closing = 3
        assert result.slide_count >= 3

    def test_extra_slide_stats(self, tmp_path: Path):
        try:
            import pptx  # noqa: F401
        except ImportError:
            pytest.skip("python-pptx not installed")

        from dazzle.pitch.generators.pptx_gen import generate_pptx

        spec = PitchSpec(
            extra_slides=[
                ExtraSlide(
                    title="Metrics",
                    layout=ExtraSlideLayout.STATS,
                    items=["100|Users", "50K|Revenue"],
                ),
            ]
        )
        ctx = PitchContext(spec=spec)
        output = tmp_path / "stats.pptx"
        result = generate_pptx(ctx, output)
        assert result.success

    def test_slide_order_respected(self, tmp_path: Path):
        try:
            import pptx  # noqa: F401
        except ImportError:
            pytest.skip("python-pptx not installed")

        from dazzle.pitch.generators.pptx_gen import generate_pptx

        spec = PitchSpec(
            problem=ProblemSpec(headline="Problem"),
            slide_order=["title", "closing"],
        )
        ctx = PitchContext(spec=spec)
        output = tmp_path / "ordered.pptx"
        result = generate_pptx(ctx, output)
        assert result.success
        # Only title + closing, problem skipped
        assert result.slide_count == 2

    def test_speaker_notes_override(self, tmp_path: Path):
        try:
            import pptx  # noqa: F401
        except ImportError:
            pytest.skip("python-pptx not installed")

        from pptx import Presentation

        from dazzle.pitch.generators.pptx_gen import generate_pptx

        spec = PitchSpec(
            problem=ProblemSpec(
                headline="Big Problem",
                speaker_notes="Custom notes here",
            ),
        )
        ctx = PitchContext(spec=spec)
        output = tmp_path / "notes.pptx"
        result = generate_pptx(ctx, output)
        assert result.success

        prs = Presentation(str(output))
        # Find the problem slide (second slide)
        problem_slide = prs.slides[1]
        notes_text = problem_slide.notes_slide.notes_text_frame.text
        assert notes_text == "Custom notes here"

    def test_chart_paths_populated(self, tmp_path: Path):
        """Test that chart_paths dict is available on PitchContext."""
        from dazzle.pitch.extractor import PitchContext as PC
        from dazzle.pitch.ir import PitchSpec as PS

        ctx = PC(spec=PS())
        assert ctx.chart_paths == {}
        ctx.chart_paths["revenue"] = tmp_path / "revenue.png"
        assert "revenue" in ctx.chart_paths


class TestNewHelpers:
    def test_create_light_slide(self):
        try:
            import pptx  # noqa: F401
        except ImportError:
            pytest.skip("python-pptx not installed")

        from pptx import Presentation
        from pptx.util import Inches

        from dazzle.pitch.generators.pptx_gen import _create_light_slide, _resolve_colors
        from dazzle.pitch.ir import BrandColors

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        colors = _resolve_colors(BrandColors())
        slide = _create_light_slide(prs, colors)
        # Light background should use the light color
        assert slide.background.fill.fore_color.rgb == colors["light"]

    def test_add_slide_heading(self):
        try:
            import pptx  # noqa: F401
        except ImportError:
            pytest.skip("python-pptx not installed")

        from pptx import Presentation
        from pptx.util import Inches

        from dazzle.pitch.generators.pptx_gen import _add_slide_heading, _resolve_colors
        from dazzle.pitch.ir import BrandColors

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        colors = _resolve_colors(BrandColors())
        y = _add_slide_heading(slide, "Test Title", colors)
        assert y == 2.0
        # Should have text box + accent bar = 2 shapes
        assert len(slide.shapes) >= 2

    def test_add_bullet_list(self):
        try:
            import pptx  # noqa: F401
        except ImportError:
            pytest.skip("python-pptx not installed")

        from pptx import Presentation
        from pptx.util import Inches

        from dazzle.pitch.generators.pptx_gen import _add_bullet_list, _resolve_colors
        from dazzle.pitch.ir import BrandColors

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        colors = _resolve_colors(BrandColors())
        final_y = _add_bullet_list(
            slide, Inches(1), 2.0, Inches(10), ["A", "B", "C"], colors, spacing=0.5
        )
        assert final_y == pytest.approx(3.5)
        assert len(slide.shapes) == 3

    def test_add_table(self):
        try:
            import pptx  # noqa: F401
        except ImportError:
            pytest.skip("python-pptx not installed")

        from pptx import Presentation
        from pptx.util import Inches

        from dazzle.pitch.generators.pptx_gen import _add_table, _resolve_colors
        from dazzle.pitch.ir import BrandColors

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        colors = _resolve_colors(BrandColors())
        shape = _add_table(
            slide,
            Inches(1),
            Inches(2),
            Inches(10),
            ["Name", "Value"],
            [["A", "1"], ["B", "2"]],
            colors,
        )
        assert shape is not None
        table = shape.table
        assert len(table.rows) == 3  # 1 header + 2 data
        assert len(table.columns) == 2

    def test_add_callout_box(self):
        try:
            import pptx  # noqa: F401
        except ImportError:
            pytest.skip("python-pptx not installed")

        from pptx import Presentation
        from pptx.util import Inches

        from dazzle.pitch.generators.pptx_gen import _add_callout_box, _resolve_colors
        from dazzle.pitch.ir import BrandColors

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        colors = _resolve_colors(BrandColors())
        shape = _add_callout_box(slide, Inches(1), Inches(2), Inches(10), "Big statement", colors)
        assert shape is not None
        # Should have box + border + text = 3 shapes
        assert len(slide.shapes) >= 3

    def test_multi_column_card_grid(self, tmp_path: Path):
        try:
            import pptx  # noqa: F401
        except ImportError:
            pytest.skip("python-pptx not installed")

        from pptx import Presentation

        from dazzle.pitch.generators.pptx_gen import generate_pptx

        spec = PitchSpec(
            extra_slides=[
                ExtraSlide(
                    title="Cards",
                    layout=ExtraSlideLayout.CARDS,
                    items=["Card 1", "Card 2", "Card 3", "Card 4"],
                ),
            ]
        )
        ctx = PitchContext(spec=spec)
        output = tmp_path / "cards.pptx"
        result = generate_pptx(ctx, output)
        assert result.success

        # Verify grid: 4 items with max 3 cols = 2 rows
        prs = Presentation(str(output))
        # Extra slide is between title and closing
        extra_slide = prs.slides[1]
        # Each card = card shape + text shape = 2; 4 cards = 8; plus heading + bar = 10
        assert len(extra_slide.shapes) >= 10

    def test_competition_uses_table(self, tmp_path: Path):
        try:
            import pptx  # noqa: F401
        except ImportError:
            pytest.skip("python-pptx not installed")

        from pptx import Presentation

        from dazzle.pitch.generators.pptx_gen import generate_pptx

        spec = PitchSpec(
            competitors=[Competitor(name="BigCo", strength="Big", weakness="Slow")],
        )
        ctx = PitchContext(spec=spec)
        output = tmp_path / "comp.pptx"
        result = generate_pptx(ctx, output)
        assert result.success

        prs = Presentation(str(output))
        # Find competition slide
        comp_slide = None
        for s in prs.slides:
            for shape in s.shapes:
                if shape.has_table:
                    comp_slide = s
                    break
        assert comp_slide is not None, "Competition slide should contain a table"

    def test_team_uses_light_background(self, tmp_path: Path):
        try:
            import pptx  # noqa: F401
        except ImportError:
            pytest.skip("python-pptx not installed")

        from pptx import Presentation

        from dazzle.pitch.generators.pptx_gen import _resolve_colors, generate_pptx
        from dazzle.pitch.ir import BrandColors

        spec = PitchSpec(
            team=TeamSpec(founders=[TeamMember(name="Alice", role="CEO")]),
        )
        ctx = PitchContext(spec=spec)
        output = tmp_path / "team.pptx"
        result = generate_pptx(ctx, output)
        assert result.success

        colors = _resolve_colors(BrandColors())
        prs = Presentation(str(output))
        # Team slide is after title
        team_slide = prs.slides[1]
        assert team_slide.background.fill.fore_color.rgb == colors["light"]

    def test_font_family_applied(self):
        try:
            import pptx  # noqa: F401
        except ImportError:
            pytest.skip("python-pptx not installed")

        from pptx import Presentation
        from pptx.util import Inches

        from dazzle.pitch.generators.pptx_gen import _add_text_box

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        txbox = _add_text_box(
            slide,
            Inches(1),
            Inches(1),
            Inches(5),
            Inches(1),
            "Hello",
            font_name="Arial",
        )
        assert txbox.text_frame.paragraphs[0].font.name == "Arial"

    def test_table_extra_slide_layout(self, tmp_path: Path):
        try:
            import pptx  # noqa: F401
        except ImportError:
            pytest.skip("python-pptx not installed")

        from pptx import Presentation

        from dazzle.pitch.generators.pptx_gen import generate_pptx

        spec = PitchSpec(
            extra_slides=[
                ExtraSlide(
                    title="Data Table",
                    layout=ExtraSlideLayout.TABLE,
                    items=["Name|Score", "Alice|95", "Bob|87"],
                ),
            ]
        )
        ctx = PitchContext(spec=spec)
        output = tmp_path / "table.pptx"
        result = generate_pptx(ctx, output)
        assert result.success

        prs = Presentation(str(output))
        extra_slide = prs.slides[1]
        has_table = any(shape.has_table for shape in extra_slide.shapes)
        assert has_table

    def test_callout_extra_slide_layout(self, tmp_path: Path):
        try:
            import pptx  # noqa: F401
        except ImportError:
            pytest.skip("python-pptx not installed")

        from dazzle.pitch.generators.pptx_gen import generate_pptx

        spec = PitchSpec(
            extra_slides=[
                ExtraSlide(
                    title="Key Insight",
                    layout=ExtraSlideLayout.CALLOUT,
                    items=["Big statement here", "Supporting detail 1", "Supporting detail 2"],
                ),
            ]
        )
        ctx = PitchContext(spec=spec)
        output = tmp_path / "callout.pptx"
        result = generate_pptx(ctx, output)
        assert result.success
        assert result.slide_count >= 3


class TestNewCosmetics:
    def test_add_divider(self):
        try:
            import pptx  # noqa: F401
        except ImportError:
            pytest.skip("python-pptx not installed")

        from pptx import Presentation

        from dazzle.pitch.generators.pptx_gen import _add_divider, _resolve_colors
        from dazzle.pitch.ir import BrandColors

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        colors = _resolve_colors(BrandColors())
        bar = _add_divider(slide, 3.0, colors)
        assert bar is not None
        assert len(slide.shapes) == 1
        # Bar should have accent color fill
        assert bar.fill.fore_color.rgb == colors["accent"]

    def test_logo_embedding_in_title_slide(self, tmp_path: Path):
        try:
            import pptx  # noqa: F401
        except ImportError:
            pytest.skip("python-pptx not installed")

        # Create a minimal valid PNG (1x1 pixel)
        import struct
        import zlib

        from pptx import Presentation

        from dazzle.pitch.generators.pptx_gen import generate_pptx

        def _make_png(path: Path) -> None:
            sig = b"\x89PNG\r\n\x1a\n"
            ihdr_data = struct.pack(">IIBBBBB", 1, 1, 8, 2, 0, 0, 0)
            ihdr_crc = zlib.crc32(b"IHDR" + ihdr_data) & 0xFFFFFFFF
            ihdr = struct.pack(">I", 13) + b"IHDR" + ihdr_data + struct.pack(">I", ihdr_crc)
            raw = b"\x00\x00\x00\x00"  # filter byte + 1 RGB pixel
            compressed = zlib.compress(raw)
            idat_crc = zlib.crc32(b"IDAT" + compressed) & 0xFFFFFFFF
            idat = (
                struct.pack(">I", len(compressed))
                + b"IDAT"
                + compressed
                + struct.pack(">I", idat_crc)
            )
            iend_crc = zlib.crc32(b"IEND") & 0xFFFFFFFF
            iend = struct.pack(">I", 0) + b"IEND" + struct.pack(">I", iend_crc)
            path.write_bytes(sig + ihdr + idat + iend)

        logo_path = tmp_path / "logo.png"
        _make_png(logo_path)

        spec = PitchSpec(
            company=CompanySpec(
                name="LogoCo",
                stage=FundingStage.SEED,
                logo_path=str(logo_path),
            ),
        )
        ctx = PitchContext(spec=spec)
        output = tmp_path / "logo_deck.pptx"
        result = generate_pptx(ctx, output)
        assert result.success

        prs = Presentation(str(output))
        title_slide = prs.slides[0]
        # Should have picture shape for logo
        has_picture = any(
            hasattr(shape, "image") and shape.image is not None
            for shape in title_slide.shapes
            if shape.shape_type == 13  # MSO_SHAPE_TYPE.PICTURE
        )
        assert has_picture, "Title slide should contain embedded logo"

    def test_pricing_tier_has_card_shapes(self, tmp_path: Path):
        try:
            import pptx  # noqa: F401
        except ImportError:
            pytest.skip("python-pptx not installed")

        from pptx import Presentation

        from dazzle.pitch.generators.pptx_gen import generate_pptx

        spec = PitchSpec(
            business_model=BusinessModelSpec(
                tiers=[
                    PricingTier(name="Free", price=0),
                    PricingTier(name="Pro", price=49, highlighted=True),
                ]
            ),
        )
        ctx = PitchContext(spec=spec)
        output = tmp_path / "pricing.pptx"
        result = generate_pptx(ctx, output)
        assert result.success

        prs = Presentation(str(output))
        # Business model slide is after title
        bm_slide = prs.slides[1]
        # Should have auto-shapes (cards) with type 5 (rounded rect)
        auto_shapes = [s for s in bm_slide.shapes if hasattr(s, "shape_type") and s.shape_type == 1]
        assert len(auto_shapes) >= 2, "Each pricing tier should have a card shape"

    def test_platform_stat_box_cards(self, tmp_path: Path):
        try:
            import pptx  # noqa: F401
        except ImportError:
            pytest.skip("python-pptx not installed")

        from pptx import Presentation

        from dazzle.pitch.generators.pptx_gen import generate_pptx

        spec = PitchSpec()
        ctx = PitchContext(
            spec=spec,
            entities=["Task", "Project"],
            surfaces=["task_list"],
        )
        output = tmp_path / "platform.pptx"
        result = generate_pptx(ctx, output)
        assert result.success

        prs = Presentation(str(output))
        # Platform slide is after title
        platform_slide = prs.slides[1]
        # Should have auto-shapes (cards) with type 5 (rounded rect)
        auto_shapes = [
            s for s in platform_slide.shapes if hasattr(s, "shape_type") and s.shape_type == 1
        ]
        assert len(auto_shapes) >= 2, "Each platform metric should have a card background"

    def test_milestones_has_dividers(self, tmp_path: Path):
        try:
            import pptx  # noqa: F401
        except ImportError:
            pytest.skip("python-pptx not installed")

        from pptx import Presentation

        from dazzle.pitch.generators.pptx_gen import generate_pptx

        spec = PitchSpec(
            milestones=MilestonesSpec(
                completed=["MVP done"],
                next_12_months=["Launch v1"],
                long_term=["Global expansion"],
            ),
        )
        ctx = PitchContext(spec=spec)
        output = tmp_path / "milestones.pptx"
        result = generate_pptx(ctx, output)
        assert result.success

        prs = Presentation(str(output))
        # Milestones slide is after title
        ms_slide = prs.slides[1]
        # Count shapes - with 3 sections and dividers between them, should have more shapes
        # than without dividers. At minimum we expect 2 divider rectangles.
        shape_count = len(ms_slide.shapes)
        # heading(text+bar) + completed(header+1 bullet) + divider + next(header+1 bullet)
        # + divider + long_term(header+1 bullet) = at least 10 shapes
        assert shape_count >= 10, f"Expected >=10 shapes with dividers, got {shape_count}"


class TestNarrativeGenerator:
    def test_generate_narrative(self, tmp_path: Path):
        from dazzle.pitch.generators.narrative import generate_narrative

        ctx = _make_full_context()
        output = tmp_path / "narrative.md"
        result = generate_narrative(ctx, output)

        assert result.success
        assert output.exists()
        content = output.read_text()
        assert "TestCo" in content
        assert "Testing is hard" in content
        assert "Market Opportunity" in content

    def test_generate_narrative_minimal(self, tmp_path: Path):
        from dazzle.pitch.generators.narrative import generate_narrative

        spec = PitchSpec()
        ctx = PitchContext(spec=spec)
        output = tmp_path / "minimal.md"
        result = generate_narrative(ctx, output)

        assert result.success
        assert output.exists()
